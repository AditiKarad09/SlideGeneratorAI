# # # # # # #
# # # # # # # import json
# # # # # # # import os
# # # # # # # import pandas as pd
# # # # # # # from bs4 import BeautifulSoup
# # # # # # # from PyPDF2 import PdfReader
# # # # # # #
# # # # # # # # Import the PDF-processing function from the separate file
# # # # # # # from processing_pdf import process_pdf
# # # # # # #
# # # # # # #
# # # # # # # def process_text(text):
# # # # # # #     """Convert plain text into structured JSON"""
# # # # # # #     return {"type": "text", "content": text.strip()}
# # # # # # #
# # # # # # #
# # # # # # # def process_json(json_str):
# # # # # # #     """Parse JSON string and return a dictionary"""
# # # # # # #     try:
# # # # # # #         return json.loads(json_str)
# # # # # # #     except json.JSONDecodeError:
# # # # # # #         return {"error": "Invalid JSON format"}
# # # # # # #
# # # # # # #
# # # # # # # def process_table(file_path):
# # # # # # #     """Read tabular data (CSV/Excel) and convert it to JSON"""
# # # # # # #     try:
# # # # # # #         if file_path.endswith(".csv"):
# # # # # # #             df = pd.read_csv(file_path)
# # # # # # #         elif file_path.endswith(".xlsx"):
# # # # # # #             df = pd.read_excel(file_path)
# # # # # # #         else:
# # # # # # #             return {"error": "Unsupported table format"}
# # # # # # #
# # # # # # #         return {"type": "table", "content": df.to_dict(orient="records")}
# # # # # # #     except Exception as e:
# # # # # # #         return {"error": f"Failed to process table: {str(e)}"}
# # # # # # #
# # # # # # #
# # # # # # # def process_css(css_str):
# # # # # # #     """Parse CSS and return a JSON representation"""
# # # # # # #     try:
# # # # # # #         soup = BeautifulSoup(css_str, "html.parser")
# # # # # # #         rules = []
# # # # # # #         for rule in soup.text.split("}"):
# # # # # # #             if "{" in rule:
# # # # # # #                 selector, properties = rule.split("{", 1)
# # # # # # #                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
# # # # # # #
# # # # # # #         return {"type": "css", "content": rules}
# # # # # # #     except Exception as e:
# # # # # # #         return {"error": f"Failed to process CSS: {str(e)}"}
# # # # # # #
# # # # # # #
# # # # # # # def validate_and_convert(file_path):
# # # # # # #     """Validate and convert different file types to JSON (or LaTeX, etc.)"""
# # # # # # #     _, ext = os.path.splitext(file_path)
# # # # # # #     ext = ext.lower()
# # # # # # #
# # # # # # #     try:
# # # # # # #         if ext == ".txt":
# # # # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # # # #                 return process_text(f.read())
# # # # # # #
# # # # # # #         elif ext == ".json":
# # # # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # # # #                 return process_json(f.read())
# # # # # # #
# # # # # # #         elif ext in [".csv", ".xlsx"]:
# # # # # # #             return process_table(file_path)
# # # # # # #
# # # # # # #         elif ext == ".css":
# # # # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # # # #                 return process_css(f.read())
# # # # # # #
# # # # # # #         elif ext == ".pdf":
# # # # # # #             # Call the PDF processing logic from processing_pdf.py
# # # # # # #             return process_pdf(file_path)
# # # # # # #
# # # # # # #         else:
# # # # # # #             return {"error": "Unsupported file format"}
# # # # # # #
# # # # # # #     except Exception as e:
# # # # # # #         return {"error": f"Failed to process file: {str(e)}"}
# # # # # # #
# # # # # # #
# # # # # # # if __name__ == "__main__":
# # # # # # #     file_path = input("Enter the file path: ")
# # # # # # #     result = validate_and_convert(file_path)
# # # # # # #     print(json.dumps(result, indent=4))
# # # # # # #
# # # # # # #########################################################################
# # # # # #
# # # # # # import json
# # # # # # import os
# # # # # # import pandas as pd
# # # # # # from bs4 import BeautifulSoup
# # # # # # from PyPDF2 import PdfReader
# # # # # #
# # # # # # # Import the PDF-processing function from the separate file
# # # # # # from processing_pdf import process_pdf
# # # # # #
# # # # # # # -------------- LLM + Plotting Imports --------------
# # # # # # import openai
# # # # # # import matplotlib.pyplot as plt
# # # # # # import re
# # # # # #
# # # # # # # Set your OpenAI API key
# # # # # # openai.api_key = os.getenv("OPENAI_API_KEY", "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")
# # # # # #
# # # # # #
# # # # # # def process_text(text):
# # # # # #     """Convert plain text into structured JSON"""
# # # # # #     return {"type": "text", "content": text.strip()}
# # # # # #
# # # # # #
# # # # # # def process_json(json_str):
# # # # # #     """Parse JSON string and return a dictionary"""
# # # # # #     try:
# # # # # #         return json.loads(json_str)
# # # # # #     except json.JSONDecodeError:
# # # # # #         return {"error": "Invalid JSON format"}
# # # # # #
# # # # # #
# # # # # # def process_table(file_path):
# # # # # #     """Read tabular data (CSV/Excel) and convert it to JSON"""
# # # # # #     try:
# # # # # #         if file_path.endswith(".csv"):
# # # # # #             df = pd.read_csv(file_path)
# # # # # #         elif file_path.endswith(".xlsx"):
# # # # # #             df = pd.read_excel(file_path)
# # # # # #         else:
# # # # # #             return {"error": "Unsupported table format"}
# # # # # #
# # # # # #         return {"type": "table", "content": df.to_dict(orient="records")}
# # # # # #     except Exception as e:
# # # # # #         return {"error": f"Failed to process table: {str(e)}"}
# # # # # #
# # # # # #
# # # # # # def process_css(css_str):
# # # # # #     """Parse CSS and return a JSON representation"""
# # # # # #     try:
# # # # # #         soup = BeautifulSoup(css_str, "html.parser")
# # # # # #         rules = []
# # # # # #         for rule in soup.text.split("}"):
# # # # # #             if "{" in rule:
# # # # # #                 selector, properties = rule.split("{", 1)
# # # # # #                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
# # # # # #
# # # # # #         return {"type": "css", "content": rules}
# # # # # #     except Exception as e:
# # # # # #         return {"error": f"Failed to process CSS: {str(e)}"}
# # # # # #
# # # # # #
# # # # # # def validate_and_convert(file_path):
# # # # # #     """Validate and convert different file types to JSON (or LaTeX, etc.)"""
# # # # # #     _, ext = os.path.splitext(file_path)
# # # # # #     ext = ext.lower()
# # # # # #
# # # # # #     try:
# # # # # #         if ext == ".txt":
# # # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # # #                 return process_text(f.read())
# # # # # #
# # # # # #         elif ext == ".json":
# # # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # # #                 return process_json(f.read())
# # # # # #
# # # # # #         elif ext in [".csv", ".xlsx"]:
# # # # # #             return process_table(file_path)
# # # # # #
# # # # # #         elif ext == ".css":
# # # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # # #                 return process_css(f.read())
# # # # # #
# # # # # #         elif ext == ".pdf":
# # # # # #             # 1) Process PDF -> returns dict with 'type' and 'content' (the LaTeX)
# # # # # #             pdf_result = process_pdf(file_path)
# # # # # #
# # # # # #             # 2) If pdf_result["type"] == "latex", parse the LaTeX for any tables
# # # # # #             if pdf_result["type"] == "latex":
# # # # # #                 latex_code = pdf_result["content"]
# # # # # #                 # Attempt to parse any tabular data and generate charts
# # # # # #                 parse_latex_and_generate_charts(latex_code)
# # # # # #
# # # # # #             # Return the PDF result as usual
# # # # # #             return pdf_result
# # # # # #
# # # # # #         else:
# # # # # #             return {"error": "Unsupported file format"}
# # # # # #
# # # # # #     except Exception as e:
# # # # # #         return {"error": f"Failed to process file: {str(e)}"}
# # # # # #
# # # # # #
# # # # # # # ------------- PARSE TABLES FROM LATEX & GENERATE CHARTS -------------
# # # # # # def parse_latex_tabular(latex_string):
# # # # # #     """
# # # # # #     Parse LaTeX string to extract table data from any \begin{tabular} ... \end{tabular}.
# # # # # #     Returns a list of tables. Each table is a list of rows (each row is a list of strings).
# # # # # #     """
# # # # # #     # 1) Normalize double backslashes to single backslash for regex parsing
# # # # # #     #    This step ensures "\begin{tabular}" is recognized if the string has "\\begin{tabular}".
# # # # # #     normalized_str = latex_string
# # # # # #     normalized_str = latex_string.replace("\\\\", "\\")
# # # # # #     normalized_str = latex_string.replace("\\", '')
# # # # # #     # print("djbdhcdbckjdsbckbscbkjsdbckjsdbckjbdcbkjdbcjks")
# # # # # #     # print(repr(normalized_str))
# # # # # #     # print("-------------------------------------")
# # # # # #     # print(normalized_str)
# # # # # #     # 2) Find everything between \begin{tabular} and \end{tabular}
# # # # # #     tabular_pattern = r"begin\{tabular\}\{.*?\}(.*?)end\{tabular\}"
# # # # # #     rows_pattern = r"(.*?)(?:\\hline|\n)"
# # # # # #
# # # # # #     # Find all tabulars in the content
# # # # # #     tabular_matches = re.findall(tabular_pattern, normalized_str, re.DOTALL)
# # # # # #
# # # # # #     # Process each tabular content
# # # # # #     dataframes = []
# # # # # #     for tabular in tabular_matches:
# # # # # #         rows = [row for row in tabular.split("\n") if row.strip() and not row.strip().startswith("hline")]
# # # # # #         table_data = [row.split("&") for row in rows]
# # # # # #         df = pd.DataFrame(table_data[1:], columns=[col.strip() for col in table_data[0]])
# # # # # #         dataframes.append(df)
# # # # # #
# # # # # #     # for d in dataframes:
# # # # # #     #     print("here------------------>")
# # # # # #     #     print(d)
# # # # # #
# # # # # #     return dataframes
# # # # # #
# # # # # #
# # # # # # def llm_recommend_chart_type(table_df, purpose="visualize data"):
# # # # # #     """
# # # # # #     Use an LLM (OpenAI GPT) to recommend the best chart type for the given table DataFrame.
# # # # # #     We'll send column names, sample rows, etc. to the LLM.
# # # # # #     """
# # # # # #     # Summarize columns and sample data
# # # # # #     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
# # # # # #     sample_records = table_df.head(3).to_dict(orient="records")
# # # # # #
# # # # # #     prompt = f"""
# # # # # # I have a table with columns: {columns_info}.
# # # # # # Sample rows: {sample_records}.
# # # # # #
# # # # # # I want to {purpose}. Which chart type is best (e.g. bar, line, scatter, etc.) and why?
# # # # # # Provide a short answer.
# # # # # # """
# # # # # #     try:
# # # # # #         response = openai.ChatCompletion.create(
# # # # # #             model="gpt-3.5-turbo",
# # # # # #             messages=[{"role": "user", "content": prompt}],
# # # # # #             temperature=0.3,
# # # # # #             max_tokens=200,
# # # # # #         )
# # # # # #         answer = response["choices"][0]["message"]["content"].strip()
# # # # # #         return answer
# # # # # #     except Exception as e:
# # # # # #         return f"Error calling OpenAI API: {str(e)}"
# # # # # #
# # # # # #
# # # # # # def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
# # # # # #     """
# # # # # #     Given a DataFrame and an LLM chart suggestion,
# # # # # #     generate a simple matplotlib plot (bar, line, scatter, hist, etc.).
# # # # # #     """
# # # # # #     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
# # # # # #     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
# # # # # #     suggestion_lower = chart_type_suggestion.lower()
# # # # # #
# # # # # #     if "line" in suggestion_lower:
# # # # # #         if len(numeric_cols) >= 2:
# # # # # #             x_col, y_col = numeric_cols[0], numeric_cols[1]
# # # # # #         elif len(numeric_cols) == 1 and len(non_numeric_cols) >= 1:
# # # # # #             x_col, y_col = non_numeric_cols[0], numeric_cols[0]
# # # # # #         else:
# # # # # #             # fallback
# # # # # #             x_col, y_col = table_df.columns[0], table_df.columns[1]
# # # # # #
# # # # # #         plt.figure()
# # # # # #         plt.plot(table_df[x_col], table_df[y_col], marker='o')
# # # # # #         plt.xlabel(str(x_col))
# # # # # #         plt.ylabel(str(y_col))
# # # # # #         plt.title(f"Line Chart: {y_col} vs {x_col}")
# # # # # #
# # # # # #     elif "scatter" in suggestion_lower:
# # # # # #         if len(numeric_cols) >= 2:
# # # # # #             x_col, y_col = numeric_cols[0], numeric_cols[1]
# # # # # #         else:
# # # # # #             x_col, y_col = table_df.columns[0], table_df.columns[1]
# # # # # #
# # # # # #         plt.figure()
# # # # # #         plt.scatter(table_df[x_col], table_df[y_col], c="blue")
# # # # # #         plt.xlabel(str(x_col))
# # # # # #         plt.ylabel(str(y_col))
# # # # # #         plt.title(f"Scatter Plot: {y_col} vs {x_col}")
# # # # # #
# # # # # #     elif "hist" in suggestion_lower:
# # # # # #         if len(numeric_cols) >= 1:
# # # # # #             col = numeric_cols[0]
# # # # # #         else:
# # # # # #             col = table_df.columns[0]
# # # # # #         plt.figure()
# # # # # #         plt.hist(table_df[col].dropna(), bins=10, color="green")
# # # # # #         plt.title(f"Histogram of {col}")
# # # # # #
# # # # # #     else:
# # # # # #         # default to bar chart
# # # # # #         if len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
# # # # # #             x_col, y_col = non_numeric_cols[0], numeric_cols[0]
# # # # # #             plt.figure()
# # # # # #             plt.bar(table_df[x_col], table_df[y_col], color="orange")
# # # # # #             plt.xlabel(str(x_col))
# # # # # #             plt.ylabel(str(y_col))
# # # # # #             plt.title(f"Bar Chart: {y_col} by {x_col}")
# # # # # #         else:
# # # # # #             # fallback
# # # # # #             col = table_df.columns[0]
# # # # # #             plt.figure()
# # # # # #             plt.bar(range(len(table_df)), pd.to_numeric(table_df[col], errors='coerce'))
# # # # # #             plt.title(f"Bar Chart of {col}")
# # # # # #
# # # # # #     plt.tight_layout()
# # # # # #     plt.savefig(output_name, dpi=150)
# # # # # #     plt.close()
# # # # # #
# # # # # #
# # # # # # def parse_latex_and_generate_charts(latex_code):
# # # # # #     """
# # # # # #     1) Parse LaTeX for tables
# # # # # #     2) For each table, build a DataFrame
# # # # # #     3) Ask LLM for recommended chart type
# # # # # #     4) Generate & save chart
# # # # # #     """
# # # # # #     # Extract tables from the LaTeX code
# # # # # #     tables_data = parse_latex_tabular(latex_code)  # Assume this function returns a list of DataFrames
# # # # # #     print(tables_data)
# # # # # #     if not tables_data:
# # # # # #         print("No tabular data found in the LaTeX output.")
# # # # # #         return
# # # # # #
# # # # # #     for i, df in enumerate(tables_data, start=1):
# # # # # #         print(f"\n--- Table {i} Detected ---")
# # # # # #         if df.empty:
# # # # # #             print("DataFrame is empty. Skipping.")
# # # # # #             continue
# # # # # #
# # # # # #         try:
# # # # # #             # Convert numeric columns where possible
# # # # # #             for col in df.columns:
# # # # # #                 df[col] = pd.to_numeric(df[col], errors="ignore")
# # # # # #
# # # # # #             # Call LLM to recommend chart type
# # # # # #             recommendation = llm_recommend_chart_type(df, purpose="visualize this table effectively")
# # # # # #             print(f"LLM Suggestion: {recommendation}")
# # # # # #
# # # # # #             # Generate the chart
# # # # # #             chart_filename = f"table_{i}_chart.png"
# # # # # #             plot_data_automatically(df, recommendation, output_name=chart_filename)
# # # # # #             print(f"Chart saved as: {chart_filename}")
# # # # # #
# # # # # #         except Exception as e:
# # # # # #             print(f"Error processing table {i}: {e}")
# # # # # #
# # # # # #
# # # # # #
# # # # # #
# # # # # # if __name__ == "__main__":
# # # # # #     file_path = input("Enter the file path: ")
# # # # # #     result = validate_and_convert(file_path)
# # # # # #     print(json.dumps(result, indent=4))
# # # # #
# # # # #
# # # # #
# # # # #
# # # # #
# # # # #
# # # # #
# # # # # import json
# # # # # import os
# # # # # import pandas as pd
# # # # # from bs4 import BeautifulSoup
# # # # # from PyPDF2 import PdfReader
# # # # # import openai
# # # # # import matplotlib.pyplot as plt
# # # # # import re
# # # # # from processing_pdf import process_pdf
# # # # #
# # # # # # Set your OpenAI API key
# # # # # openai.api_key = os.getenv("OPENAI_API_KEY")
# # # # #
# # # # # # Functions for processing different file types
# # # # # def process_text(text):
# # # # #     """Convert plain text into structured JSON"""
# # # # #     return {"type": "text", "content": text.strip()}
# # # # #
# # # # # def process_json(json_str):
# # # # #     """Parse JSON string and return a dictionary"""
# # # # #     try:
# # # # #         return json.loads(json_str)
# # # # #     except json.JSONDecodeError:
# # # # #         return {"error": "Invalid JSON format"}
# # # # #
# # # # # def process_table(file_path):
# # # # #     """Read tabular data (CSV/Excel) and convert it to JSON"""
# # # # #     try:
# # # # #         if file_path.endswith(".csv"):
# # # # #             df = pd.read_csv(file_path)
# # # # #         elif file_path.endswith(".xlsx"):
# # # # #             df = pd.read_excel(file_path)
# # # # #         else:
# # # # #             return {"error": "Unsupported table format"}
# # # # #         return {"type": "table", "content": df.to_dict(orient="records")}
# # # # #     except Exception as e:
# # # # #         return {"error": f"Failed to process table: {str(e)}"}
# # # # #
# # # # # def process_css(css_str):
# # # # #     """Parse CSS and return a JSON representation"""
# # # # #     try:
# # # # #         soup = BeautifulSoup(css_str, "html.parser")
# # # # #         rules = []
# # # # #         for rule in soup.text.split("}"):
# # # # #             if "{" in rule:
# # # # #                 selector, properties = rule.split("{", 1)
# # # # #                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
# # # # #         return {"type": "css", "content": rules}
# # # # #     except Exception as e:
# # # # #         return {"error": f"Failed to process CSS: {str(e)}"}
# # # # #
# # # # # def validate_and_convert(file_path):
# # # # #     """Validate and convert different file types to JSON (or LaTeX, etc.)"""
# # # # #     _, ext = os.path.splitext(file_path)
# # # # #     ext = ext.lower()
# # # # #
# # # # #     try:
# # # # #         if ext == ".txt":
# # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # #                 return process_text(f.read())
# # # # #         elif ext == ".json":
# # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # #                 return process_json(f.read())
# # # # #         elif ext in [".csv", ".xlsx"]:
# # # # #             return process_table(file_path)
# # # # #         elif ext == ".css":
# # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # #                 return process_css(f.read())
# # # # #         elif ext == ".pdf":
# # # # #             pdf_result = process_pdf(file_path)  # Assuming you have a process_pdf function
# # # # #             if pdf_result["type"] == "latex":
# # # # #                 latex_code = pdf_result["content"]
# # # # #                 parse_latex_and_generate_charts(latex_code)
# # # # #             return pdf_result
# # # # #         else:
# # # # #             return {"error": "Unsupported file format"}
# # # # #     except Exception as e:
# # # # #         return {"error": f"Failed to process file: {str(e)}"}
# # # # #
# # # # # # Parse LaTeX for tables
# # # # # def parse_latex_tabular(latex_string):
# # # # #     normalized_str = latex_string.replace("\\\\", "\\").replace("\\", "")
# # # # #     tabular_pattern = r"begin\{tabular\}\{.*?\}(.*?)end\{tabular\}"
# # # # #     tabular_matches = re.findall(tabular_pattern, normalized_str, re.DOTALL)
# # # # #     dataframes = []
# # # # #     for tabular in tabular_matches:
# # # # #         rows = [row for row in tabular.split("\n") if row.strip() and not row.strip().startswith("hline")]
# # # # #         table_data = [row.split("&") for row in rows]
# # # # #         df = pd.DataFrame(table_data[1:], columns=[col.strip() for col in table_data[0]])
# # # # #         dataframes.append(df)
# # # # #     return dataframes
# # # # #
# # # # # # Call OpenAI to recommend a chart type
# # # # # def llm_recommend_chart_type(table_df, purpose="visualize data"):
# # # # #     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
# # # # #     sample_records = table_df.head(3).to_dict(orient="records")
# # # # #     prompt = f"""
# # # # # I have a table with columns: {columns_info}.
# # # # # Sample rows: {sample_records}.
# # # # # I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, etc.) and why?
# # # # # Provide a short answer.
# # # # # """
# # # # #     try:
# # # # #         response = openai.ChatCompletion.create(
# # # # #             model="gpt-4",
# # # # #             messages=[
# # # # #                 {"role": "system", "content": "You are a data visualization expert."},
# # # # #                 {"role": "user", "content": prompt},
# # # # #             ],
# # # # #             temperature=0.3,
# # # # #             max_tokens=200,
# # # # #         )
# # # # #         return response["choices"][0]["message"]["content"].strip()
# # # # #     except Exception as e:
# # # # #         return f"Error calling OpenAI API: {str(e)}"
# # # # #
# # # # # # Plot data automatically
# # # # # def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
# # # # #     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
# # # # #     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
# # # # #     suggestion_lower = chart_type_suggestion.lower()
# # # # #
# # # # #     plt.figure()
# # # # #     if "line" in suggestion_lower and len(numeric_cols) >= 2:
# # # # #         plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
# # # # #         plt.title("Line Chart")
# # # # #     elif "scatter" in suggestion_lower and len(numeric_cols) >= 2:
# # # # #         plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
# # # # #         plt.title("Scatter Plot")
# # # # #     elif "hist" in suggestion_lower and len(numeric_cols) >= 1:
# # # # #         plt.hist(table_df[numeric_cols[0]].dropna(), bins=10, color="green")
# # # # #         plt.title("Histogram")
# # # # #     elif len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
# # # # #         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
# # # # #         plt.title("Bar Chart")
# # # # #     else:
# # # # #         plt.bar(range(len(table_df)), range(len(table_df)))
# # # # #         plt.title("Default Bar Chart")
# # # # #     plt.savefig(output_name)
# # # # #     plt.close()
# # # # #
# # # # # # Parse LaTeX and generate charts
# # # # # def parse_latex_and_generate_charts(latex_code):
# # # # #     tables_data = parse_latex_tabular(latex_code)
# # # # #     if not tables_data:
# # # # #         print("No tabular data found.")
# # # # #         return
# # # # #     for i, df in enumerate(tables_data, start=1):
# # # # #         print(f"\n--- Table {i} Detected ---")
# # # # #         if df.empty:
# # # # #             print("DataFrame is empty. Skipping.")
# # # # #             continue
# # # # #         try:
# # # # #             for col in df.columns:
# # # # #                 df[col] = pd.to_numeric(df[col], errors="ignore")
# # # # #             recommendation = llm_recommend_chart_type(df, purpose="visualize this table effectively")
# # # # #             print(f"LLM Suggestion: {recommendation}")
# # # # #             chart_filename = f"table_{i}_chart.png"
# # # # #             plot_data_automatically(df, recommendation, output_name=chart_filename)
# # # # #             print(f"Chart saved as: {chart_filename}")
# # # # #         except Exception as e:
# # # # #             print(f"Error processing table {i}: {e}")
# # # # #
# # # # # if __name__ == "__main__":
# # # # #     file_path = input("Enter the file path: ")
# # # # #     result = validate_and_convert(file_path)
# # # # #     print(json.dumps(result, indent=4))
# # # #
# # # #
# # # #
# # # # # import json
# # # # # import os
# # # # # import pandas as pd
# # # # # from bs4 import BeautifulSoup
# # # # # from PyPDF2 import PdfReader
# # # # # import google.generativeai as genai
# # # # # import matplotlib.pyplot as plt
# # # # # import re
# # # # # from processing_pdf import process_pdf
# # # # #
# # # # #
# # # # # # Configure Gemini API
# # # # # genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")
# # # # # model = genai.GenerativeModel("gemini-1.5-flash")
# # # # #
# # # # # # Functions for processing different file types
# # # # # def process_text(text):
# # # # #     """Convert plain text into structured JSON"""
# # # # #     return {"type": "text", "content": text.strip()}
# # # # #
# # # # # def process_json(json_str):
# # # # #     """Parse JSON string and return a dictionary"""
# # # # #     try:
# # # # #         return json.loads(json_str)
# # # # #     except json.JSONDecodeError:
# # # # #         return {"error": "Invalid JSON format"}
# # # # #
# # # # # def process_table(file_path):
# # # # #     """Read tabular data (CSV/Excel) and convert it to JSON"""
# # # # #     try:
# # # # #         if file_path.endswith(".csv"):
# # # # #             df = pd.read_csv(file_path)
# # # # #         elif file_path.endswith(".xlsx"):
# # # # #             df = pd.read_excel(file_path)
# # # # #         else:
# # # # #             return {"error": "Unsupported table format"}
# # # # #         return {"type": "table", "content": df.to_dict(orient="records")}
# # # # #     except Exception as e:
# # # # #         return {"error": f"Failed to process table: {str(e)}"}
# # # # #
# # # # # def process_css(css_str):
# # # # #     """Parse CSS and return a JSON representation"""
# # # # #     try:
# # # # #         soup = BeautifulSoup(css_str, "html.parser")
# # # # #         rules = []
# # # # #         for rule in soup.text.split("}"):
# # # # #             if "{" in rule:
# # # # #                 selector, properties = rule.split("{", 1)
# # # # #                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
# # # # #         return {"type": "css", "content": rules}
# # # # #     except Exception as e:
# # # # #         return {"error": f"Failed to process CSS: {str(e)}"}
# # # # #
# # # # # def validate_and_convert(file_path):
# # # # #     """Validate and convert different file types to JSON (or LaTeX, etc.)"""
# # # # #     _, ext = os.path.splitext(file_path)
# # # # #     ext = ext.lower()
# # # # #
# # # # #     try:
# # # # #         if ext == ".txt":
# # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # #                 return process_text(f.read())
# # # # #         elif ext == ".json":
# # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # #                 return process_json(f.read())
# # # # #         elif ext in [".csv", ".xlsx"]:
# # # # #             return process_table(file_path)
# # # # #         elif ext == ".css":
# # # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # # #                 return process_css(f.read())
# # # # #         elif ext == ".pdf":
# # # # #             pdf_result = process_pdf(file_path)  # Assuming you have a process_pdf function
# # # # #             if pdf_result["type"] == "latex":
# # # # #                 latex_code = pdf_result["content"]
# # # # #                 parse_latex_and_generate_charts(latex_code)
# # # # #             return pdf_result
# # # # #         else:
# # # # #             return {"error": "Unsupported file format"}
# # # # #     except Exception as e:
# # # # #         return {"error": f"Failed to process file: {str(e)}"}
# # # # #
# # # # # # Parse LaTeX for tables
# # # # # def parse_latex_tabular(latex_string):
# # # # #     normalized_str = latex_string.replace("\\\\", "\\").replace("\\", "")
# # # # #     tabular_pattern = r"begin\{tabular\}\{.*?\}(.*?)end\{tabular\}"
# # # # #     tabular_matches = re.findall(tabular_pattern, normalized_str, re.DOTALL)
# # # # #     dataframes = []
# # # # #     for tabular in tabular_matches:
# # # # #         rows = [row for row in tabular.split("\n") if row.strip() and not row.strip().startswith("hline")]
# # # # #         table_data = [row.split("&") for row in rows]
# # # # #         df = pd.DataFrame(table_data[1:], columns=[col.strip() for col in table_data[0]])
# # # # #         dataframes.append(df)
# # # # #     return dataframes
# # # # #
# # # # # # Call Gemini to recommend a chart type
# # # # # def gemini_recommend_chart_type(table_df, purpose="visualize data"):
# # # # #     """
# # # # #     Uses Gemini to recommend a chart type for the given table.
# # # # #
# # # # #     Args:
# # # # #         table_df: pandas DataFrame containing the table data.
# # # # #         purpose: Optional purpose for the visualization (default: "visualize data").
# # # # #
# # # # #     Returns:
# # # # #         A string containing the recommended chart type.
# # # # #     """
# # # # #     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
# # # # #     sample_records = table_df.head(3).to_dict(orient="records")
# # # # #
# # # # #     prompt = f"""
# # # # # I have a table with columns: {columns_info}.
# # # # # Sample rows: {sample_records}.
# # # # # I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, etc.)?
# # # # # Provide a short and concise answer.
# # # # # """
# # # # #
# # # # #     try:
# # # # #         response = model.generate_content(prompt)
# # # # #         return response.text.strip()
# # # # #     except Exception as e:
# # # # #         return f"Error calling Gemini API: {str(e)}"
# # # # #
# # # # # # Plot data automatically
# # # # # def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
# # # # #     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
# # # # #     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
# # # # #     suggestion_lower = chart_type_suggestion.lower()
# # # # #
# # # # #     plt.figure()
# # # # #     if "line" in suggestion_lower and len(numeric_cols) >= 2:
# # # # #         plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
# # # # #         plt.title("Line Chart")
# # # # #     elif "scatter" in suggestion_lower and len(numeric_cols) >= 2:
# # # # #         plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
# # # # #         plt.title("Scatter Plot")
# # # # #     elif "hist" in suggestion_lower and len(numeric_cols) >= 1:
# # # # #         plt.hist(table_df[numeric_cols[0]].dropna(), bins=10, color="green")
# # # # #         plt.title("Histogram")
# # # # #     elif len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
# # # # #         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
# # # # #         plt.title("Bar Chart")
# # # # #     else:
# # # # #         plt.bar(range(len(table_df)), range(len(table_df)))
# # # # #         plt.title("Default Bar Chart")
# # # # #     plt.savefig(output_name)
# # # # #     plt.close()
# # # # #
# # # # # # Parse LaTeX and generate charts
# # # # # def parse_latex_and_generate_charts(latex_code):
# # # # #     tables_data = parse_latex_tabular(latex_code)
# # # # #     if not tables_data:
# # # # #         print("No tabular data found.")
# # # # #         return
# # # # #     for i, df in enumerate(tables_data, start=1):
# # # # #         print(f"\n--- Table {i} Detected ---")
# # # # #         if df.empty:
# # # # #             print("DataFrame is empty. Skipping.")
# # # # #             continue
# # # # #         try:
# # # # #             for col in df.columns:
# # # # #                 df[col] = pd.to_numeric(df[col], errors="ignore")
# # # # #             recommendation = gemini_recommend_chart_type(df, purpose="visualize this table effectively")
# # # # #             print(f"Gemini Suggestion: {recommendation}")
# # # # #             chart_filename = f"table_{i}_chart.png"
# # # # #             plot_data_automatically(df, recommendation, output_name=chart_filename)
# # # # #             print(f"Chart saved as: {chart_filename}")
# # # # #         except Exception as e:
# # # # #             print(f"Error processing table {i}: {e}")
# # # # #
# # # # # if __name__ == "__main__":
# # # # #     file_path = input("Enter the file path: ")
# # # # #     result = validate_and_convert(file_path)
# # # # #     print(json.dumps(result, indent=4))
# # # #
# # # #
# # # #
# # # #
# # # # import json
# # # # import os
# # # # import pandas as pd
# # # # from bs4 import BeautifulSoup
# # # # import seaborn as sns
# # # # from PyPDF2 import PdfReader
# # # # import google.generativeai as genai
# # # # import matplotlib.pyplot as plt
# # # # import re
# # # # from processing_pdf import process_pdf  # Assuming you have a process_pdf function
# # # # import pptx
# # # # from pptx.util import Inches
# # # #
# # # # # Configure Gemini API
# # # # genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")
# # # # model = genai.GenerativeModel("gemini-1.5-flash")
# # # #
# # # # # Functions for processing different file types
# # # # def process_text(text):
# # # #     """Convert plain text into structured JSON"""
# # # #     return {"type": "text", "content": text.strip()}
# # # #
# # # # def process_json(json_str):
# # # #     """Parse JSON string and return a dictionary"""
# # # #     try:
# # # #         return json.loads(json_str)
# # # #     except json.JSONDecodeError:
# # # #         return {"error": "Invalid JSON format"}
# # # #
# # # # def process_table(file_path):
# # # #     """Read tabular data (CSV/Excel) and convert it to JSON"""
# # # #     try:
# # # #         if file_path.endswith(".csv"):
# # # #             df = pd.read_csv(file_path)
# # # #         elif file_path.endswith(".xlsx"):
# # # #             df = pd.read_excel(file_path)
# # # #         else:
# # # #             return {"error": "Unsupported table format"}
# # # #         return {"type": "table", "content": df.to_dict(orient="records")}
# # # #     except Exception as e:
# # # #         return {"error": f"Failed to process table: {str(e)}"}
# # # #
# # # # def process_css(css_str):
# # # #     """Parse CSS and return a JSON representation"""
# # # #     try:
# # # #         soup = BeautifulSoup(css_str, "html.parser")
# # # #         rules = []
# # # #         for rule in soup.text.split("}"):
# # # #             if "{" in rule:
# # # #                 selector, properties = rule.split("{", 1)
# # # #                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
# # # #         return {"type": "css", "content": rules}
# # # #     except Exception as e:
# # # #         return {"error": f"Failed to process CSS: {str(e)}"}
# # # #
# # # # def validate_and_convert(file_path):
# # # #     """Validate and convert different file types to JSON (or LaTeX, etc.)"""
# # # #     _, ext = os.path.splitext(file_path)
# # # #     ext = ext.lower()
# # # #
# # # #     try:
# # # #         if ext == ".txt":
# # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # #                 return process_text(f.read())
# # # #         elif ext == ".json":
# # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # #                 return process_json(f.read())
# # # #         elif ext in [".csv", ".xlsx"]:
# # # #             return process_table(file_path)
# # # #         elif ext == ".css":
# # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # #                 return process_css(f.read())
# # # #         elif ext == ".pdf":
# # # #             pdf_result = process_pdf(file_path)  # Assuming you have a process_pdf function
# # # #             if pdf_result["type"] == "latex":
# # # #                 return {"type": "latex", "content": pdf_result["content"]}
# # # #             else:
# # # #                 return pdf_result
# # # #         else:
# # # #             return {"error": "Unsupported file format"}
# # # #     except Exception as e:
# # # #         return {"error": f"Failed to process file: {str(e)}"}
# # # #
# # # # # Parse LaTeX for tables
# # # # def parse_latex_tabular(latex_string):
# # # #     normalized_str = latex_string.replace("\\\\", "\\").replace("\\", "")
# # # #     tabular_pattern = r"begin\{tabular\}\{.*?\}(.*?)end\{tabular\}"
# # # #     tabular_matches = re.findall(tabular_pattern, normalized_str, re.DOTALL)
# # # #     dataframes = []
# # # #     for tabular in tabular_matches:
# # # #         rows = [row for row in tabular.split("\n") if row.strip() and not row.strip().startswith("hline")]
# # # #         table_data = [row.split("&") for row in rows]
# # # #         df = pd.DataFrame(table_data[1:], columns=[col.strip() for col in table_data[0]])
# # # #         dataframes.append(df)
# # # #     return dataframes
# # # #
# # # # # Call Gemini to recommend a chart type
# # # # def gemini_recommend_chart_type(table_df, purpose="visualize data"):
# # # #     """
# # # #     Uses Gemini to recommend a chart type for the given table.
# # # #
# # # #     Args:
# # # #         table_df: pandas DataFrame containing the table data.
# # # #         purpose: Optional purpose for the visualization (default: "visualize data").
# # # #
# # # #     Returns:
# # # #         A string containing the recommended chart type.
# # # #     """
# # # #     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
# # # #     sample_records = table_df.head(3).to_dict(orient="records")
# # # #
# # # #     prompt = f"""
# # # # I have a table with columns: {columns_info}.
# # # # Sample rows: {sample_records}.
# # # # I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, pie etc.)?
# # # # Provide a short and concise answer.
# # # # """
# # # #
# # # #     try:
# # # #         response = model.generate_content(prompt)
# # # #         return response.text.strip()
# # # #     except Exception as e:
# # # #         return f"Error calling Gemini API: {str(e)}"
# # # #
# # # # # Plot data automatically
# # # # # def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
# # # # #     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
# # # # #     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
# # # # #     suggestion_lower = chart_type_suggestion.lower()
# # # # #
# # # # #     plt.figure()
# # # # #     if "line" in suggestion_lower and len(numeric_cols) >= 2:
# # # # #         plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
# # # # #         plt.title("Line Chart")
# # # # #     elif "scatter" in suggestion_lower and len(numeric_cols) >= 2:
# # # # #         plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
# # # # #         plt.title("Scatter Plot")
# # # # #     elif "hist" in suggestion_lower and len(numeric_cols) >= 1:
# # # # #         plt.hist(table_df[numeric_cols[0]].dropna(), bins=10, color="green")
# # # # #         plt.title("Histogram")
# # # # #     elif len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
# # # # #         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
# # # # #         plt.title("Bar Chart")
# # # # #     else:
# # # # #         plt.bar(range(len(table_df)), range(len(table_df)))
# # # # #         plt.title("Default Bar Chart")
# # # # #     plt.savefig(output_name)
# # # # #     plt.close()
# # # # #
# # # # # def generate_pptx_from_gemini(data, title="Presentation"):
# # # # #     """Generates a PowerPoint presentation based on Gemini's instructions.
# # # # #
# # # # #     Args:
# # # # #         data: A dictionary containing the processed data (e.g., from validate_and_convert).
# # # # #         title: The title of the presentation (default: "Presentation").
# # # # #
# # # # #     Returns:
# # # # #         None (saves the presentation to "output.pptx")
# # # # #     """
# # # # #     prs = pptx.Presentation()
# # # # #
# # # # #     # Title slide
# # # # #     title_slide_layout = prs.slide_layouts[0]
# # # # #     slide = prs.slides.add_slide(title_slide_layout)
# # # # #     title_shape = slide.shapes.title
# # # # #     title_shape.text = title
# # # # #
# # # # #     # Generate Gemini prompt for presentation structure
# # # # #     prompt = f"""
# # # # #     Create a PowerPoint presentation based on the following information:
# # # # #
# # # # #     **Data:** {json.dumps(data, indent=4)}
# # # # #
# # # # #     **Instructions:**
# # # # #     - **Structure:** Suggest a logical slide order and content for each slide.
# # # # #     - **Visualizations:** Recommend appropriate charts and graphs for relevant data.
# # # # #     - **Text Formatting:** Suggest headings, bullet points, and concise text for each slide.
# # # # #
# # # # #     **Example Output:**
# # # # #     - **Slide 1 (Title):** Presentation Title
# # # # #       - **Content:** Brief introduction or overview.
# # # # #     - **Slide 2 (Data Summary):**
# # # # #       - **Content:** Summary of key findings or data points.
# # # # #       - **Chart:** Bar chart showing [description of chart].
# # # # #     - **Slide 3 (Analysis):**
# # # # #       - **Content:** Analysis and insights based on the data.
# # # # #
# # # # #     Provide the output in a similar structured format.
# # # # #     """
# # # # #
# # # # #     try:
# # # # #         response = model.generate_content(prompt)
# # # # #         print(response.text)  # Print Gemini's suggestions for analysis and debugging
# # # # #
# # # # #         # Parse Gemini's response to create slide data
# # # # #         slides_data = []
# # # # #         for slide_info in response.text.strip().split("- "):
# # # # #             if slide_info:
# # # # #                 slide_data = {}
# # # # #                 try:
# # # # #                     slide_title, slide_content_and_chart = slide_info.split(":")
# # # # #                     slide_data["title"] = slide_title.strip()
# # # # #                     slide_content = slide_content_and_chart.split("- Chart:")[0].strip()
# # # # #                     slide_data["content"] = slide_content
# # # # #
# # # # #                     if "Chart:" in slide_content_and_chart:
# # # # #                         chart_type = slide_content_and_chart.split("- Chart:")[1].strip()
# # # # #                         # Generate chart based on chart_type (you'll need to implement this logic)
# # # # #                         # ... (Chart generation logic based on chart_type and data)
# # # # #                         chart_filename = f"chart_{len(slides_data)}.png"
# # # # #                         # ... (Use plot_data_automatically or a more advanced charting function)
# # # # #                         slide_data["chart_path"] = chart_filename
# # # # #                 except ValueError:
# # # # #                     print(f"Error parsing slide information: {slide_info}")
# # # # #
# # # # #                 slides_data.append(slide_data)
# # # # #
# # # # #         # Create slides in the presentation
# # # # #         for slide_data in slides_data:
# # # # #             slide_layout = prs.slide_layouts[1]  # Use a content layout (e.g., title + content)
# # # # #             slide = prs.slides.add_slide(slide_layout)
# # # # #
# # # # #             # Add title
# # # # #             title_shape = slide.shapes.title
# # # # #             title_shape.text = slide_data.get("title", "")
# # # # #
# # # # #             # Add content (text, bullet points)
# # # # #             body_shape = slide.placeholders[1]
# # # # #             if "content" in slide_data:
# # # # #                 body_shape.text = slide_data["content"]
# # # # #
# # # # #             # Add chart (if available)
# # # # #             if "chart_path" in slide_data:
# # # # #                 left = Inches(1)
# # # # #                 top = Inches(2)
# # # # #                 width = Inches(6)
# # # # #                 height = Inches(4)
# # # # #                 pic = slide.shapes.add_picture(slide_data["chart_path"], left, top, width, height)
# # # # #
# # # # #         prs.save("output.pptx")
# # # # #
# # # # #     except Exception as e:
# # # # #         print(f"Error calling Gemini API or creating slides: {str(e)}")
# # # # #
# # # # # if __name__ == "__main__":
# # # # #     file_path = input("Enter the file path: ")
# # # # #     data = validate_and_convert(file_path)
# # # # #     generate_pptx_from_gemini(data)
# # # #
# # # # def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
# # # #     """
# # # #     Generates a chart based on the table data and chart type suggestion.
# # # #
# # # #     Args:
# # # #         table_df: pandas DataFrame containing the data.
# # # #         chart_type_suggestion: String suggesting the chart type (e.g., "bar", "line").
# # # #         output_name: Optional filename to save the chart image (default: "chart.png").
# # # #     """
# # # #
# # # #     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
# # # #     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
# # # #     suggestion_lower = chart_type_suggestion.lower()
# # # #
# # # #     plt.figure()
# # # #     if suggestion_lower in ["line", "scatter"] and len(numeric_cols) >= 2:
# # # #         if suggestion_lower == "line":
# # # #             plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
# # # #         else:
# # # #             plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
# # # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # # #     elif suggestion_lower in ["hist", "distribution"] and len(numeric_cols) >= 1:
# # # #         sns.histplot(table_df[numeric_cols[0]])
# # # #         plt.title(f"{chart_type_suggestion.capitalize()} Plot")
# # # #     elif suggestion_lower in ["bar", "column"] and len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
# # # #         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
# # # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # # #     elif suggestion_lower in ["bar", "column"] and len(non_numeric_cols) >= 1:
# # # #         sns.countplot(x=table_df[non_numeric_cols[0]])
# # # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # # #     else:
# # # #         print(f"Unsupported chart type: {chart_type_suggestion}")
# # # #         return
# # # #
# # # #     plt.savefig(output_name)
# # # #     plt.close()
# # # #
# # # # def generate_pptx_from_gemini(data, title="Presentation"):
# # # #     """Generates a PowerPoint presentation based on Gemini's instructions.
# # # #
# # # #     Args:
# # # #         data: A dictionary containing the processed data (e.g., from validate_and_convert).
# # # #         title: The title of the presentation (default: "Presentation").
# # # #
# # # #     Returns:
# # # #         None (saves the presentation to "output.pptx")
# # # #     """
# # # #     prs = pptx.Presentation()
# # # #
# # # #     # Title slide
# # # #     title_slide_layout = prs.slide_layouts[0]
# # # #     slide = prs.slides.add_slide(title_slide_layout)
# # # #     title_shape = slide.shapes.title
# # # #     title_shape.text = title
# # # #
# # # #     # Generate Gemini prompt for presentation structure
# # # #     prompt = f"""
# # # #     Create a PowerPoint presentation based on the following information:
# # # #
# # # #     **Data:** {json.dumps(data, indent=4)}
# # # #
# # # #     **Instructions:**
# # # #     - **Structure:** Suggest a logical slide order and content for each slide.
# # # #     - **Visualizations:** Recommend appropriate charts and graphs for relevant data.
# # # #     - **Text Formatting:** Suggest headings, bullet points, and concise text for each slide.
# # # #
# # # #     **Example Output:**
# # # #     - **Slide 1 (Title):** Presentation Title
# # # #       - **Content:** Brief introduction or overview.
# # # #     - **Slide 2 (Data Summary):**
# # # #       - **Content:** Summary of key findings or data points.
# # # #       - **Chart:** Bar chart showing [description of chart].
# # # #     - **Slide 3 (Analysis):**
# # # #       - **Content:** Analysis and insights based on the data.
# # # #
# # # #     Provide the output in a similar structured format.
# # # #     """
# # # #
# # # #     try:
# # # #         response = model.generate_content(prompt)
# # # #         print(response.text)  # Print Gemini's suggestions for analysis and debugging
# # # #
# # # #         # Parse Gemini's response to create slide data
# # # #         slides_data = []
# # # #         for slide_info in response.text.strip().split("- "):
# # # #             if slide_info:
# # # #                 slide_data = {}
# # # #                 try:
# # # #                     slide_title, slide_content_and_chart = slide_info.split(":")
# # # #                     slide_data["title"] = slide_title.strip()
# # # #                     slide_content = slide_content_and_chart.split("- Chart:")[0].strip()
# # # #                     slide_data["content"] = slide_content
# # # #
# # # #                     if "Chart:" in slide_content_and_chart:
# # # #                         chart_type = slide_content_and_chart.split("- Chart:")[1].strip()
# # # #                         chart_filename = f"chart_{len(slides_data)}.png"
# # # #                         slide_data["chart_path"] = chart_filename
# # # #                 except ValueError:
# # # #                     print(f"Error parsing slide information: {slide_info}")
# # # #                     slide_data["title"] = slide_info
# # # #                     slide_data["content"] = ""
# # # #
# # # #                 slides_data.append(slide_data)
# # # #
# # # #         # Create slides in the presentation
# # # #         for slide_data in slides_data:
# # # #             slide_layout = prs.slide_layouts[1]  # Use a content layout (e.g., title + content)
# # # #             slide = prs.slides.add_slide(slide_layout)
# # # #
# # # #             # Add title
# # # #             title_shape = slide.shapes.title
# # # #             title_shape.text = slide_data.get("title", "")
# # # #
# # # #             # Add content (text, bullet points)
# # # #             body_shape = slide.placeholders[1]
# # # #             if "content" in slide_data:
# # # #                 body_shape.text = slide_data["content"]
# # # #
# # # #             # Add chart (if available)
# # # #             if "chart_path" in slide_data:
# # # #                 left = Inches(1)
# # # #                 top = Inches(2)
# # # #                 width = Inches(6)
# # # #                 height = Inches(4)
# # # #                 pic = slide.shapes.add_picture(slide_data["chart_path"], left, top, width, height)
# # # #
# # # #         prs.save("output.pptx")
# # # #
# # # #     except Exception as e:
# # # #         print(f"Error calling Gemini API or creating slides: {str(e)}")
# # # #
# # # # if __name__ == "__main__":
# # # #     file_path = input("Enter the file path: ")
# # # #     data = validate_and_convert(file_path)
# # # #     generate_pptx_from_gemini(data)
# # #
# # #
# # # # import json
# # # # import os
# # # # import pandas as pd
# # # # from bs4 import BeautifulSoup
# # # # import seaborn as sns
# # # # from PyPDF2 import PdfReader
# # # # import google.generativeai as genai
# # # # import matplotlib.pyplot as plt
# # # # import re
# # # # from processing_pdf import process_pdf  # Assuming you have a process_pdf function
# # # # import pptx
# # # # from pptx.util import Inches
# # # #
# # # # # Configure Gemini API
# # # # genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")
# # # # model = genai.GenerativeModel("gemini-1.5-flash")
# # # #
# # # # # --------------------------------------------------------------------
# # # # # 1. FILE PROCESSING FUNCTIONS
# # # # # --------------------------------------------------------------------
# # # #
# # # # def process_text(text):
# # # #     """Convert plain text into structured JSON."""
# # # #     return {"type": "text", "content": text.strip()}
# # # #
# # # # def process_json(json_str):
# # # #     """Parse JSON string and return a dictionary."""
# # # #     try:
# # # #         return json.loads(json_str)
# # # #     except json.JSONDecodeError:
# # # #         return {"error": "Invalid JSON format"}
# # # #
# # # # def process_table(file_path):
# # # #     """Read tabular data (CSV/Excel) and convert it to JSON."""
# # # #     try:
# # # #         if file_path.endswith(".csv"):
# # # #             df = pd.read_csv(file_path)
# # # #         elif file_path.endswith(".xlsx"):
# # # #             df = pd.read_excel(file_path)
# # # #         else:
# # # #             return {"error": "Unsupported table format"}
# # # #         return {"type": "table", "content": df.to_dict(orient="records")}
# # # #     except Exception as e:
# # # #         return {"error": f"Failed to process table: {str(e)}"}
# # # #
# # # # def process_css(css_str):
# # # #     """Parse CSS and return a JSON representation."""
# # # #     try:
# # # #         soup = BeautifulSoup(css_str, "html.parser")
# # # #         rules = []
# # # #         for rule in soup.text.split("}"):
# # # #             if "{" in rule:
# # # #                 selector, properties = rule.split("{", 1)
# # # #                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
# # # #         return {"type": "css", "content": rules}
# # # #     except Exception as e:
# # # #         return {"error": f"Failed to process CSS: {str(e)}"}
# # # #
# # # # def validate_and_convert(file_path):
# # # #     """Validate and convert different file types to a dictionary representation."""
# # # #     _, ext = os.path.splitext(file_path)
# # # #     ext = ext.lower()
# # # #
# # # #     try:
# # # #         if ext == ".txt":
# # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # #                 return process_text(f.read())
# # # #         elif ext == ".json":
# # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # #                 return process_json(f.read())
# # # #         elif ext in [".csv", ".xlsx"]:
# # # #             return process_table(file_path)
# # # #         elif ext == ".css":
# # # #             with open(file_path, "r", encoding="utf-8") as f:
# # # #                 return process_css(f.read())
# # # #         elif ext == ".pdf":
# # # #             pdf_result = process_pdf(file_path)  # Assuming you have a process_pdf function
# # # #             if pdf_result.get("type") == "latex":
# # # #                 return {"type": "latex", "content": pdf_result["content"]}
# # # #             else:
# # # #                 return pdf_result
# # # #         else:
# # # #             return {"error": "Unsupported file format"}
# # # #     except Exception as e:
# # # #         return {"error": f"Failed to process file: {str(e)}"}
# # # #
# # # # # --------------------------------------------------------------------
# # # # # 2. HELPER FUNCTIONS FOR LATEX (OPTIONAL)
# # # # # --------------------------------------------------------------------
# # # #
# # # # def parse_latex_tabular(latex_string):
# # # #     """
# # # #     Extracts tabular environments from LaTeX and returns a list of DataFrames.
# # # #     """
# # # #     normalized_str = latex_string.replace("\\\\", "\\").replace("\\", "")
# # # #     tabular_pattern = r"begin\{tabular\}\{.*?\}(.*?)end\{tabular\}"
# # # #     tabular_matches = re.findall(tabular_pattern, normalized_str, re.DOTALL)
# # # #     dataframes = []
# # # #     for tabular in tabular_matches:
# # # #         rows = [
# # # #             row for row in tabular.split("\n")
# # # #             if row.strip() and not row.strip().startswith("hline")
# # # #         ]
# # # #         table_data = [row.split("&") for row in rows]
# # # #         df = pd.DataFrame(table_data[1:], columns=[col.strip() for col in table_data[0]])
# # # #         dataframes.append(df)
# # # #     return dataframes
# # # #
# # # # # --------------------------------------------------------------------
# # # # # 3. GEMINI + CHART FUNCTIONS
# # # # # --------------------------------------------------------------------
# # # #
# # # # def gemini_recommend_chart_type(table_df, purpose="visualize data"):
# # # #     """
# # # #     Calls Gemini to recommend a chart type for the given table.
# # # #     """
# # # #     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
# # # #     sample_records = table_df.head(3).to_dict(orient="records")
# # # #
# # # #     prompt = f"""
# # # # I have a table with columns: {columns_info}.
# # # # Sample rows: {sample_records}.
# # # # I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, pie etc.)?
# # # # Provide a short and concise answer.
# # # # """
# # # #
# # # #     try:
# # # #         response = model.generate_content(prompt)
# # # #         return response.text.strip()
# # # #     except Exception as e:
# # # #         return f"Error calling Gemini API: {str(e)}"
# # # #
# # # # def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
# # # #     """
# # # #     Generates and saves a chart based on the table data and chart type suggestion.
# # # #     """
# # # #     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
# # # #     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
# # # #     suggestion_lower = chart_type_suggestion.lower()
# # # #
# # # #     plt.figure()
# # # #     if suggestion_lower in ["line", "scatter"] and len(numeric_cols) >= 2:
# # # #         if suggestion_lower == "line":
# # # #             plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
# # # #         else:
# # # #             plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
# # # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # # #     elif suggestion_lower in ["hist", "distribution"] and len(numeric_cols) >= 1:
# # # #         sns.histplot(table_df[numeric_cols[0]])
# # # #         plt.title(f"{chart_type_suggestion.capitalize()} Plot")
# # # #     elif suggestion_lower in ["bar", "column"] and len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
# # # #         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
# # # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # # #     elif suggestion_lower in ["bar", "column"] and len(non_numeric_cols) >= 1:
# # # #         sns.countplot(x=table_df[non_numeric_cols[0]])
# # # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # # #     else:
# # # #         print(f"Unsupported or ambiguous chart type: {chart_type_suggestion}")
# # # #         return
# # # #
# # # #     plt.tight_layout()
# # # #     plt.savefig(output_name)
# # # #     plt.close()
# # # #
# # # # # --------------------------------------------------------------------
# # # # # 4. PPT GENERATION FUNCTION
# # # # # --------------------------------------------------------------------
# # # # def generate_pptx_from_gemini(data, slides_title="Presentation"):
# # # #     """
# # # #     Generates a PowerPoint based on Gemini's instructions for the given data.
# # # #     1) Calls Gemini for a proposed slide structure.
# # # #     2) Parses the response using a robust regex.
# # # #     3) Creates slides, including any charts that were generated.
# # # #     4) Saves the final PPTX.
# # # #     """
# # # #     prs = pptx.Presentation()
# # # #
# # # #     # --- 1. Title Slide ---
# # # #     title_slide_layout = prs.slide_layouts[0]
# # # #     slide = prs.slides.add_slide(title_slide_layout)
# # # #     slide.shapes.title.text = slides_title
# # # #
# # # #     # --- 2. Build Gemini Prompt for PPT Structure ---
# # # #     gemini_prompt = f"""
# # # # Create a PowerPoint presentation based on the following information:
# # # #
# # # # **Data:** {json.dumps(data, indent=4)}
# # # #
# # # # **Instructions:**
# # # # - **Structure:** Suggest a logical slide order and content for each slide.
# # # # - **Visualizations:** Mention any chart or table if needed ("- Chart:" or "- Table: ...").
# # # # - **Text Formatting:** Suggest headings, bullet points, and concise text for each slide.
# # # #
# # # # Provide the output in a bullet form like:
# # # # - **Slide 1 (Title): Some Title**
# # # #   - Content: ...
# # # #   - Bullet Points: ...
# # # #   - Chart: ...
# # # # ...
# # # #     """
# # # #
# # # #     # --- 3. Call Gemini to get slide structure ---
# # # #     try:
# # # #         response = model.generate_content(gemini_prompt)
# # # #         gemini_text = response.text.strip()
# # # #         print("[Gemini Slide Structure Suggestion]:")
# # # #         print(gemini_text)
# # # #     except Exception as e:
# # # #         print(f"Error calling Gemini API: {str(e)}")
# # # #         return
# # # #
# # # #     # --- 4. Regex to Capture Each Slide ---
# # # #     # Regex Explanation:
# # # #     # - We look for lines starting with: `- **Slide `
# # # #     # - Capture the slide number, e.g., `1`
# # # #     # - Capture the text in parentheses, e.g. `(Title)`
# # # #     # - Then capture everything until the next `- **Slide` or the end of the string.
# # # #     slide_pattern = (
# # # #         r"- \*\*Slide\s+(\d+)\s*\(([^)]+)\):"   # e.g. "- **Slide 1 (Title):"
# # # #         r"(.*?)"                               # Slide body (non-greedy)
# # # #         r"(?=- \*\*Slide|\Z)"                  # until we hit the next slide or end
# # # #     )
# # # #     matches = re.findall(slide_pattern, gemini_text, flags=re.DOTALL)
# # # #
# # # #     slides_data = []
# # # #     for match in matches:
# # # #         slide_number = match[0].strip()
# # # #         slide_subtitle = match[1].strip()
# # # #         slide_body = match[2].strip()
# # # #
# # # #         # Now inside slide_body we can look for lines like:
# # # #         #   "- Content: ..."
# # # #         #   "- Bullet Points: ..."
# # # #         #   "- Chart: ..."
# # # #         #   "- Table: ..."
# # # #         content_match = re.search(r"- Content:\s*(.*?)(?=\n- |$)", slide_body, flags=re.DOTALL)
# # # #         bullet_match = re.search(r"- Bullet Points?:\s*(.*?)(?=\n- |$)", slide_body, flags=re.DOTALL)
# # # #         chart_match = re.search(r"- Chart:\s*(.*?)(?=\n- |$)", slide_body, flags=re.DOTALL)
# # # #         table_match = re.search(r"- Table:\s*(.*?)(?=\n- |$)", slide_body, flags=re.DOTALL)
# # # #
# # # #         # Clean up the text if found
# # # #         content_text = content_match.group(1).strip() if content_match else ""
# # # #         bullet_text = bullet_match.group(1).strip() if bullet_match else ""
# # # #         chart_text = chart_match.group(1).strip() if chart_match else ""
# # # #         table_text = table_match.group(1).strip() if table_match else ""
# # # #
# # # #         # Store in a dict
# # # #         slide_dict = {
# # # #             "slide_number": slide_number,
# # # #             "slide_subtitle": slide_subtitle,
# # # #             "content": content_text,
# # # #             "bullets": bullet_text.split("\n") if bullet_text else [],
# # # #             "chart": chart_text,
# # # #             "table": table_text
# # # #         }
# # # #         slides_data.append(slide_dict)
# # # #
# # # #     # --- 5. Create Slides in PPTX ---
# # # #     for sdata in slides_data:
# # # #         slide_layout = prs.slide_layouts[1]  # Title + Content layout
# # # #         slide = prs.slides.add_slide(slide_layout)
# # # #
# # # #         # Title is "Slide {num} ({subtitle})" or something
# # # #         slide_title = f"Slide {sdata['slide_number']} ({sdata['slide_subtitle']})"
# # # #         slide.shapes.title.text = slide_title
# # # #
# # # #         # Content shape
# # # #         body_shape = slide.placeholders[1]
# # # #         text_lines = []
# # # #
# # # #         # Add main content
# # # #         if sdata["content"]:
# # # #             text_lines.append(sdata["content"])
# # # #
# # # #         # Add bullet points
# # # #         if sdata["bullets"]:
# # # #             bullet_block = "\n".join(f"• {b.strip()}" for b in sdata["bullets"] if b.strip())
# # # #             text_lines.append(bullet_block)
# # # #
# # # #         # Combine them into one text
# # # #         combined_text = "\n\n".join([t for t in text_lines if t])
# # # #         body_shape.text = combined_text
# # # #
# # # #         # If there's a chart mention, and we actually generated a chart file, place it
# # # #         # e.g., data["chart_filename"] was set earlier if you generated a chart
# # # #         if sdata["chart"] and "chart_filename" in data:
# # # #             chart_path = data["chart_filename"]
# # # #             if os.path.exists(chart_path):
# # # #                 left = Inches(1)
# # # #                 top = Inches(2)
# # # #                 width = Inches(6)
# # # #                 height = Inches(4)
# # # #                 slide.shapes.add_picture(chart_path, left, top, width, height)
# # # #             else:
# # # #                 print(f"Chart file not found: {chart_path}")
# # # #
# # # #         # If there's a table mention, handle similarly...
# # # #         # This depends on how you want to place a table image or create a real PPT table.
# # # #
# # # #     # --- 6. Save the PPT ---
# # # #     output_filename = "output2.pptx"
# # # #     prs.save(output_filename)
# # # #     print(f"Presentation saved as {output_filename}.")
# # # #
# # # # # --------------------------------------------------------------------
# # # # # 5. MAIN FLOW
# # # # # --------------------------------------------------------------------
# # # #
# # # # if __name__ == "__main__":
# # # #     # 1) Ask user for file path
# # # #     file_path = input("Enter the file path: ")
# # # #
# # # #     # 2) Validate + convert file
# # # #     data = validate_and_convert(file_path)
# # # #     if "error" in data:
# # # #         print(f"Error processing file: {data['error']}")
# # # #         exit()
# # # #
# # # #     # 3) If the data is tabular, call Gemini to recommend a chart and then generate it
# # # #     if data.get("type") == "table":
# # # #         # Convert data["content"] (list of dicts) to a DataFrame
# # # #         df = pd.DataFrame(data["content"])
# # # #
# # # #         # Get chart type suggestion from Gemini
# # # #         suggested_chart_type = gemini_recommend_chart_type(df)
# # # #         print(f"Gemini suggests a '{suggested_chart_type}' chart.")
# # # #
# # # #         # Generate chart
# # # #         chart_filename = "chart_0.png"
# # # #         plot_data_automatically(df, suggested_chart_type, output_name=chart_filename)
# # # #
# # # #         # Store the chart filename back into data so we can access it in PPT generation
# # # #         data["chart_filename"] = chart_filename
# # # #         data["chart_type"] = suggested_chart_type
# # # #
# # # #     # 4) Generate a PPT with Gemini instructions (this also references the chart if it exists)
# # # #     generate_pptx_from_gemini(data, slides_title="My Gemini Presentation")
# # #
# # #
# # #
# # #
# # # import json
# # # import os
# # # import pandas as pd
# # # from bs4 import BeautifulSoup
# # # import seaborn as sns
# # # from PyPDF2 import PdfReader
# # # import google.generativeai as genai
# # # import matplotlib.pyplot as plt
# # # import re
# # # from processing_pdf import process_pdf  # Assuming you have a process_pdf function
# # # import pptx
# # # from pptx.util import Inches
# # # from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
# # #
# # # # Configure Gemini API
# # # genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")
# # # model = genai.GenerativeModel("gemini-1.5-flash")
# # #
# # # # --------------------------------------------------------------------
# # # # 1. FILE PROCESSING FUNCTIONS
# # # # --------------------------------------------------------------------
# # #
# # # def process_text(text):
# # #     """Convert plain text into structured JSON."""
# # #     return {"type": "text", "content": text.strip()}
# # #
# # # def process_json(json_str):
# # #     """Parse JSON string and return a dictionary."""
# # #     try:
# # #         return json.loads(json_str)
# # #     except json.JSONDecodeError:
# # #         return {"error": "Invalid JSON format"}
# # #
# # # def process_table(file_path):
# # #     """Read tabular data (CSV/Excel) and convert it to JSON."""
# # #     try:
# # #         if file_path.endswith(".csv"):
# # #             df = pd.read_csv(file_path)
# # #         elif file_path.endswith(".xlsx"):
# # #             df = pd.read_excel(file_path)
# # #         else:
# # #             return {"error": "Unsupported table format"}
# # #         return {"type": "table", "content": df.to_dict(orient="records")}
# # #     except Exception as e:
# # #         return {"error": f"Failed to process table: {str(e)}"}
# # #
# # # def process_css(css_str):
# # #     """Parse CSS and return a JSON representation."""
# # #     try:
# # #         soup = BeautifulSoup(css_str, "html.parser")
# # #         rules = []
# # #         for rule in soup.text.split("}"):
# # #             if "{" in rule:
# # #                 selector, properties = rule.split("{", 1)
# # #                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
# # #         return {"type": "css", "content": rules}
# # #     except Exception as e:
# # #         return {"error": f"Failed to process CSS: {str(e)}"}
# # #
# # # def validate_and_convert(file_path):
# # #     """Validate and convert different file types to a dictionary representation."""
# # #     _, ext = os.path.splitext(file_path)
# # #     ext = ext.lower()
# # #
# # #     try:
# # #         if ext == ".txt":
# # #             with open(file_path, "r", encoding="utf-8") as f:
# # #                 return process_text(f.read())
# # #         elif ext == ".json":
# # #             with open(file_path, "r", encoding="utf-8") as f:
# # #                 return process_json(f.read())
# # #         elif ext in [".csv", ".xlsx"]:
# # #             return process_table(file_path)
# # #         elif ext == ".css":
# # #             with open(file_path, "r", encoding="utf-8") as f:
# # #                 return process_css(f.read())
# # #         elif ext == ".pdf":
# # #             pdf_result = process_pdf(file_path)  # Assuming you have a process_pdf function
# # #             if pdf_result.get("type") == "latex":
# # #                 return {"type": "latex", "content": pdf_result["content"]}
# # #             else:
# # #                 return pdf_result
# # #         else:
# # #             return {"error": "Unsupported file format"}
# # #     except Exception as e:
# # #         return {"error": f"Failed to process file: {str(e)}"}
# # #
# # # # --------------------------------------------------------------------
# # # # 2. HELPER FUNCTIONS FOR LATEX (OPTIONAL)
# # # # --------------------------------------------------------------------
# # #
# # # def parse_latex_tabular(latex_string):
# # #     """
# # #     Extracts tabular environments from LaTeX and returns a list of DataFrames.
# # #     """
# # #     normalized_str = latex_string.replace("\\\\", "\\").replace("\\", "")
# # #     tabular_pattern = r"begin\{tabular\}\{.*?\}(.*?)end\{tabular\}"
# # #     tabular_matches = re.findall(tabular_pattern, normalized_str, re.DOTALL)
# # #     dataframes = []
# # #     for tabular in tabular_matches:
# # #         rows = [
# # #             row for row in tabular.split("\n")
# # #             if row.strip() and not row.strip().startswith("hline")
# # #         ]
# # #         table_data = [row.split("&") for row in rows]
# # #         df = pd.DataFrame(table_data[1:], columns=[col.strip() for col in table_data[0]])
# # #         dataframes.append(df)
# # #     return dataframes
# # #
# # # # --------------------------------------------------------------------
# # # # 3. GEMINI + CHART FUNCTIONS
# # # # --------------------------------------------------------------------
# # #
# # # def gemini_recommend_chart_type(table_df, purpose="visualize data"):
# # #     """
# # #     Calls Gemini to recommend a chart type for the given table.
# # #     """
# # #     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
# # #     sample_records = table_df.head(3).to_dict(orient="records")
# # #
# # #     prompt = f"""
# # # I have a table with columns: {columns_info}.
# # # Sample rows: {sample_records}.
# # # I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, pie etc.)?
# # # Provide a short and concise answer.
# # # """
# # #
# # #     try:
# # #         response = model.generate_content(prompt)
# # #         return response.text.strip()
# # #     except Exception as e:
# # #         return f"Error calling Gemini API: {str(e)}"
# # #
# # # def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
# # #     """
# # #     Generates and saves a chart based on the table data and chart type suggestion.
# # #     """
# # #     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
# # #     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
# # #     suggestion_lower = chart_type_suggestion.lower()
# # #
# # #     plt.figure()
# # #     if suggestion_lower in ["line", "scatter"] and len(numeric_cols) >= 2:
# # #         if suggestion_lower == "line":
# # #             plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
# # #         else:
# # #             plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
# # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # #     elif suggestion_lower in ["hist", "distribution"] and len(numeric_cols) >= 1:
# # #         sns.histplot(table_df[numeric_cols[0]])
# # #         plt.title(f"{chart_type_suggestion.capitalize()} Plot")
# # #     elif suggestion_lower in ["bar", "column"] and len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
# # #         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
# # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # #     elif suggestion_lower in ["bar", "column"] and len(non_numeric_cols) >= 1:
# # #         sns.countplot(x=table_df[non_numeric_cols[0]])
# # #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# # #     else:
# # #         print(f"Unsupported or ambiguous chart type: {chart_type_suggestion}")
# # #         return
# # #
# # #     plt.tight_layout()
# # #     plt.savefig(output_name)
# # #     plt.close()
# # #
# # # # --------------------------------------------------------------------
# # # # 4. PPT GENERATION FUNCTION
# # # # --------------------------------------------------------------------
# # # def generate_pptx_from_gemini(data, slides_title="Presentation"):
# # #     """
# # #     Generates a PowerPoint based on Gemini's instructions for the given data.
# # #     1) Calls Gemini for a proposed slide structure.
# # #     2) Parses the response using a robust regex.
# # #     3) Creates slides, including any charts that were generated.
# # #     4) Saves the final PPTX.
# # #     """
# # #     prs = pptx.Presentation()
# # #
# # #     # --- 1. Title Slide ---
# # #     title_slide_layout = prs.slide_layouts[0]
# # #     slide = prs.slides.add_slide(title_slide_layout)
# # #     slide.shapes.title.text = slides_title
# # #
# # #     # --- 2. Build Gemini Prompt for PPT Structure ---
# # #     gemini_prompt = f"""
# # # Create a PowerPoint presentation based on the following information:
# # #
# # # **Data:** {json.dumps(data, indent=4)}
# # #
# # # **Instructions:**
# # # - **Structure:** Suggest a logical slide order and content for each slide.
# # # - **Visualizations:** Mention any chart or table if needed ("- Chart:" or "- Table: ...").
# # # - **Text Formatting:** Suggest headings, bullet points, and concise text for each slide.
# # #
# # # Provide the output in a bullet form like:
# # # - **Slide 1 (Title): Some Title**
# # #   - Content: ...
# # #   - Bullet Points: ...
# # #   - Chart: ...
# # # ...
# # #     """
# # #
# # #     # --- 3. Call Gemini to get slide structure ---
# # #     try:
# # #         response = model.generate_content(gemini_prompt)
# # #         gemini_text = response.text.strip()
# # #         print("[Gemini Slide Structure Suggestion]:")
# # #         print(gemini_text)
# # #     except Exception as e:
# # #         print(f"Error calling Gemini API: {str(e)}")
# # #         return
# # #
# # #     # --- 4. Regex to Capture Each Slide ---
# # #     # Regex Explanation:
# # #     # - We look for lines starting with: `- **Slide `
# # #     # - Capture the slide number, e.g., `1`
# # #     # - Capture the text in parentheses, e.g. `(Title)`
# # #     # - Then capture everything until the next `- **Slide` or the end of the string.
# # #     slide_pattern = (
# # #         r"- \*\*Slide\s+(\d+)\s*\(([^)]+)\):"   # e.g. "- **Slide 1 (Title):"
# # #         r"(.*?)"                               # Slide body (non-greedy)
# # #         r"(?=- \*\*Slide|\Z)"                  # until we hit the next slide or end
# # #     )
# # #     matches = re.findall(slide_pattern, gemini_text, flags=re.DOTALL)
# # #
# # #     slides_data = []
# # #     for match in matches:
# # #         slide_number = match[0].strip()
# # #         slide_subtitle = match[1].strip()
# # #         slide_body = match[2].strip()
# # #
# # #         # Now inside slide_body we can look for lines like:
# # #         #   "- Content: ..."
# # #         #   "- Bullet Points: ..."
# # #         #   "- Chart: ..."
# # #         #   "- Table: ..."
# # #         content_match = re.search(r"- Content:\s*(.*?)(?=\n- |$)", slide_body, flags=re.DOTALL)
# # #         bullet_match = re.search(r"- Bullet Points?:\s*(.*?)(?=\n- |$)", slide_body, flags=re.DOTALL)
# # #         chart_match = re.search(r"- Chart:\s*(.*?)(?=\n- |$)", slide_body, flags=re.DOTALL)
# # #         table_match = re.search(r"- Table:\s*(.*?)(?=\n- |$)", slide_body, flags=re.DOTALL)
# # #
# # #         # Clean up the text if found
# # #         content_text = content_match.group(1).strip() if content_match else ""
# # #         bullet_text = bullet_match.group(1).strip() if bullet_match else ""
# # #         chart_text = chart_match.group(1).strip() if chart_match else ""
# # #         table_text = table_match.group(1).strip() if table_match else ""
# # #
# # #         # Convert bullet_text to a list of bullet lines by splitting on newlines
# # #         # Adjust to your liking if Gemini uses different separators (periods, semicolons, etc.).
# # #         bullet_lines = [b.strip() for b in bullet_text.split("\n") if b.strip()]
# # #
# # #         # Store in a dict
# # #         slide_dict = {
# # #             "slide_number": slide_number,
# # #             "slide_subtitle": slide_subtitle,
# # #             "content": content_text,
# # #             "bullets": bullet_lines,
# # #             "chart": chart_text,
# # #             "table": table_text
# # #         }
# # #         slides_data.append(slide_dict)
# # #
# # #     # --- 5. Create Slides in PPTX ---
# # #     for sdata in slides_data:
# # #         slide_layout = prs.slide_layouts[1]  # Title + Content layout
# # #         slide = prs.slides.add_slide(slide_layout)
# # #
# # #         # Title is "Slide {num} ({subtitle})"
# # #         slide_title = f"Slide {sdata['slide_number']} ({sdata['slide_subtitle']})"
# # #         slide.shapes.title.text = slide_title
# # #
# # #         # Content shape
# # #         body_shape = slide.placeholders[1]
# # #         text_frame = body_shape.text_frame
# # #         text_frame.clear()
# # #
# # #         # 5a) Add main content as a single paragraph
# # #         if sdata["content"]:
# # #             p = text_frame.add_paragraph()
# # #             p.text = sdata["content"]
# # #             p.bullet = False
# # #             p.level = 0
# # #             p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
# # #
# # #         # 5b) Add bullet points as separate paragraphs
# # #         for bullet_line in sdata["bullets"]:
# # #             p = text_frame.add_paragraph()
# # #             p.text = bullet_line
# # #             p.bullet = True
# # #             p.level = 1
# # #             p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
# # #
# # #         # 5c) If there's a chart mention, and we actually generated a chart file, place it
# # #         if sdata["chart"] and "chart_filename" in data:
# # #             chart_path = data["chart_filename"]
# # #             if os.path.exists(chart_path):
# # #                 left = Inches(1)
# # #                 top = Inches(2)
# # #                 width = Inches(6)
# # #                 height = Inches(4)
# # #                 slide.shapes.add_picture(chart_path, left, top, width, height)
# # #             else:
# # #                 print(f"Chart file not found: {chart_path}")
# # #
# # #         # 5d) If there's a table mention, handle similarly (create a table or insert an image).
# # #         # You can either generate a screenshot of the table and insert it,
# # #         # or build a PPTX table object programmatically.
# # #
# # #     # --- 6. Save the PPT ---
# # #     output_filename = "output3.pptx"
# # #     prs.save(output_filename)
# # #     print(f"Presentation saved as {output_filename}.")
# # #
# # # # --------------------------------------------------------------------
# # # # 5. MAIN FLOW
# # # # --------------------------------------------------------------------
# # #
# # # if __name__ == "__main__":
# # #     # 1) Ask user for file path
# # #     file_path = input("Enter the file path: ")
# # #
# # #     # 2) Validate + convert file
# # #     data = validate_and_convert(file_path)
# # #     if "error" in data:
# # #         print(f"Error processing file: {data['error']}")
# # #         exit()
# # #
# # #     # 3) If the data is tabular, call Gemini to recommend a chart and then generate it
# # #     if data.get("type") == "table":
# # #         # Convert data["content"] (list of dicts) to a DataFrame
# # #         df = pd.DataFrame(data["content"])
# # #
# # #         # Get chart type suggestion from Gemini
# # #         suggested_chart_type = gemini_recommend_chart_type(df)
# # #         print(f"Gemini suggests a '{suggested_chart_type}' chart.")
# # #
# # #         # Generate chart
# # #         chart_filename = "chart_0.png"
# # #         plot_data_automatically(df, suggested_chart_type, output_name=chart_filename)
# # #
# # #         # Store the chart filename back into data so we can access it in PPT generation
# # #         data["chart_filename"] = chart_filename
# # #         data["chart_type"] = suggested_chart_type
# # #
# # #     # 4) Generate a PPT with Gemini instructions (this also references the chart if it exists)
# # #     generate_pptx_from_gemini(data, slides_title="My Gemini Presentation")
# #
# # import json
# # import os
# # import re
# # import pandas as pd
# # from bs4 import BeautifulSoup
# # import seaborn as sns
# # from PyPDF2 import PdfReader
# # import google.generativeai as genai
# # import matplotlib.pyplot as plt
# # from processing_pdf import process_pdf  # Assuming you have a process_pdf function
# # import pptx
# # from pptx.util import Inches
# # from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
# #
# # # --------------------------------------------------------------------
# # # 0. CONFIGURE GEMINI
# # # --------------------------------------------------------------------
# # genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")  # Replace with your actual key
# # model = genai.GenerativeModel("gemini-1.5-flash")
# #
# #
# # # --------------------------------------------------------------------
# # # 1. FILE PROCESSING FUNCTIONS
# # # --------------------------------------------------------------------
# #
# # def process_text(text):
# #     """Convert plain text into structured JSON."""
# #     return {"type": "text", "content": text.strip()}
# #
# # def process_json(json_str):
# #     """Parse JSON string and return a dictionary."""
# #     try:
# #         return json.loads(json_str)
# #     except json.JSONDecodeError:
# #         return {"error": "Invalid JSON format"}
# #
# # def process_table(file_path):
# #     """Read tabular data (CSV/Excel) and convert it to JSON."""
# #     try:
# #         if file_path.endswith(".csv"):
# #             df = pd.read_csv(file_path)
# #         elif file_path.endswith(".xlsx"):
# #             df = pd.read_excel(file_path)
# #         else:
# #             return {"error": "Unsupported table format"}
# #         return {"type": "table", "content": df.to_dict(orient="records")}
# #     except Exception as e:
# #         return {"error": f"Failed to process table: {str(e)}"}
# #
# # def process_css(css_str):
# #     """Parse CSS and return a JSON representation."""
# #     try:
# #         soup = BeautifulSoup(css_str, "html.parser")
# #         rules = []
# #         for rule in soup.text.split("}"):
# #             if "{" in rule:
# #                 selector, properties = rule.split("{", 1)
# #                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
# #         return {"type": "css", "content": rules}
# #     except Exception as e:
# #         return {"error": f"Failed to process CSS: {str(e)}"}
# #
# # def validate_and_convert(file_path):
# #     """Validate and convert different file types to a dictionary representation."""
# #     _, ext = os.path.splitext(file_path)
# #     ext = ext.lower()
# #
# #     try:
# #         if ext == ".txt":
# #             with open(file_path, "r", encoding="utf-8") as f:
# #                 return process_text(f.read())
# #         elif ext == ".json":
# #             with open(file_path, "r", encoding="utf-8") as f:
# #                 return process_json(f.read())
# #         elif ext in [".csv", ".xlsx"]:
# #             return process_table(file_path)
# #         elif ext == ".css":
# #             with open(file_path, "r", encoding="utf-8") as f:
# #                 return process_css(f.read())
# #         elif ext == ".pdf":
# #             # Assuming you have a process_pdf function that returns {"type": "...", "content": "..."}
# #             pdf_result = process_pdf(file_path)
# #             if pdf_result.get("type") == "latex":
# #                 return {"type": "latex", "content": pdf_result["content"]}
# #             else:
# #                 return pdf_result
# #         else:
# #             return {"error": "Unsupported file format"}
# #     except Exception as e:
# #         return {"error": f"Failed to process file: {str(e)}"}
# #
# #
# # # --------------------------------------------------------------------
# # # 2. (OPTIONAL) LATEX TABLE PARSING HELPER
# # # --------------------------------------------------------------------
# #
# # def parse_latex_tabular(latex_string):
# #     """
# #     Extracts tabular environments from LaTeX and returns a list of DataFrames.
# #     """
# #     normalized_str = latex_string.replace("\\\\", "\\").replace("\\", "")
# #     tabular_pattern = r"begin\{tabular\}\{.*?\}(.*?)end\{tabular\}"
# #     tabular_matches = re.findall(tabular_pattern, normalized_str, re.DOTALL)
# #     dataframes = []
# #     for tabular in tabular_matches:
# #         rows = [
# #             row for row in tabular.split("\n")
# #             if row.strip() and not row.strip().startswith("hline")
# #         ]
# #         table_data = [row.split("&") for row in rows]
# #         df = pd.DataFrame(table_data[1:], columns=[col.strip() for col in table_data[0]])
# #         dataframes.append(df)
# #     return dataframes
# #
# #
# # # --------------------------------------------------------------------
# # # 3. GEMINI + CHART FUNCTIONS
# # # --------------------------------------------------------------------
# #
# # def gemini_recommend_chart_type(table_df, purpose="visualize data"):
# #     """
# #     Calls Gemini to recommend a chart type for the given table.
# #     """
# #     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
# #     sample_records = table_df.head(3).to_dict(orient="records")
# #
# #     prompt = f"""
# # I have a table with columns: {columns_info}.
# # Sample rows: {sample_records}.
# # I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, pie etc.)?
# # Provide a short and concise answer.
# # """
# #
# #     try:
# #         response = model.generate_content(prompt)
# #         return response.text.strip()
# #     except Exception as e:
# #         return f"Error calling Gemini API: {str(e)}"
# #
# #
# # def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
# #     """
# #     Generates and saves a chart based on the table data and chart type suggestion.
# #     """
# #     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
# #     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
# #     suggestion_lower = chart_type_suggestion.lower()
# #
# #     plt.figure()
# #     if suggestion_lower in ["line", "scatter"] and len(numeric_cols) >= 2:
# #         if suggestion_lower == "line":
# #             plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
# #         else:
# #             plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
# #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# #     elif suggestion_lower in ["hist", "distribution"] and len(numeric_cols) >= 1:
# #         sns.histplot(table_df[numeric_cols[0]])
# #         plt.title(f"{chart_type_suggestion.capitalize()} Plot")
# #     elif suggestion_lower in ["bar", "column"] and len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
# #         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
# #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# #     elif suggestion_lower in ["bar", "column"] and len(non_numeric_cols) >= 1:
# #         sns.countplot(x=table_df[non_numeric_cols[0]])
# #         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
# #     else:
# #         print(f"Unsupported or ambiguous chart type: {chart_type_suggestion}")
# #         return
# #
# #     plt.tight_layout()
# #     plt.savefig(output_name)
# #     plt.close()
# #
# #
# # # --------------------------------------------------------------------
# # # 4. HELPER: SPLIT BULLETS ON PERIOD
# # # --------------------------------------------------------------------
# #
# # def split_bullets_on_period(bullet_text):
# #     """
# #     Splits bullet_text on a period + optional space => distinct bullet lines.
# #     Example:
# #       "First point. Second point. Third point."
# #       => ["First point", "Second point", "Third point"]
# #     """
# #     bullet_text = bullet_text.strip()
# #     lines = re.split(r"\.\s+|\.$", bullet_text)
# #     lines = [l.strip() for l in lines if l.strip()]
# #     return lines
# #
# #
# # # --------------------------------------------------------------------
# # # 5. PPT GENERATION FUNCTION
# # # --------------------------------------------------------------------
# #
# # def generate_pptx_from_gemini(data, slides_title="Presentation"):
# #     """
# #     1) Calls Gemini for a proposed slide structure.
# #     2) Uses a regex to parse each slide block.
# #     3) Within each slide block, extracts lines:
# #        - Content: ...
# #        - Bullet Points: ...
# #        - Chart: ...
# #        - Table: ...
# #     4) Creates slides with distinct paragraphs for content and bullet points.
# #     5) Saves the final PPTX.
# #     """
# #     prs = pptx.Presentation()
# #
# #     # --- 1. Title Slide ---
# #     title_slide_layout = prs.slide_layouts[0]
# #     slide = prs.slides.add_slide(title_slide_layout)
# #     slide.shapes.title.text = slides_title
# #
# #     # --- 2. Build Gemini Prompt for PPT Structure ---
# #     gemini_prompt = f"""
# # Create a PowerPoint presentation based on the following information:
# #
# # **Data:** {json.dumps(data, indent=4)}
# #
# # **Instructions:**
# # - **Structure:** Suggest a logical slide order and content for each slide.
# # - **Visualizations:** Mention any chart or table if needed ("- Chart:" or "- Table: ...").
# # - **Text Formatting:** Provide headings, bullet points, and concise text for each slide.
# # - Use a format like:
# #     - **Slide 1 (Title): My Title**
# #       - Content: ...
# #       - Bullet Points: ...
# #       - Table: ...
# #       - Chart: ...
# #     - **Slide 2 (Title): Another Slide**
# #       ...
# #
# # Return only the structured outline.
# #     """
# #
# #     # --- 3. Call Gemini for the presentation structure ---
# #     try:
# #         response = model.generate_content(gemini_prompt)
# #         gemini_text = response.text.strip()
# #         print("[Gemini Slide Structure Suggestion]:")
# #         print(gemini_text)
# #     except Exception as e:
# #         print(f"Error calling Gemini API: {str(e)}")
# #         return
# #
# #     # --- 4. Parse each slide using regex ---
# #     # Looking for blocks like:
# #     #  - **Slide 1 (Title): XXX
# #     #     - Content: ...
# #     #     - Bullet Points: ...
# #     # etc.
# #     slide_pattern = (
# #         r"- \*\*Slide\s+(\d+)\s*\(([^)]+)\):"  # e.g., "- **Slide 1 (Title):"
# #         r"(.*?)"                               # capture everything
# #         r"(?=- \*\*Slide|\Z)"                  # until next slide or end
# #     )
# #     raw_slides = re.findall(slide_pattern, gemini_text, flags=re.DOTALL)
# #
# #     slides_data = []
# #     for slide_num, slide_subtitle, slide_body in raw_slides:
# #         slide_num = slide_num.strip()
# #         slide_subtitle = slide_subtitle.strip()
# #         slide_body = slide_body.strip()
# #
# #         # We'll parse line by line for:
# #         #   - Content:
# #         #   - Bullet Points:
# #         #   - Chart:
# #         #   - Table:
# #         # Everything else merges into "content".
# #         content_lines = []
# #         bullet_text_accumulator = []
# #         chart_text = ""
# #         table_text = ""
# #
# #         # Split the captured body into lines (ignore blank lines)
# #         lines = [ln.strip() for ln in slide_body.splitlines() if ln.strip()]
# #
# #         for line in lines:
# #             content_match = re.match(r"^- Content:\s*(.*)$", line, re.IGNORECASE)
# #             bullet_match = re.match(r"^- Bullet Points?:\s*(.*)$", line, re.IGNORECASE)
# #             chart_match = re.match(r"^- Chart:\s*(.*)$", line, re.IGNORECASE)
# #             table_match = re.match(r"^- Table:\s*(.*)$", line, re.IGNORECASE)
# #
# #             if content_match:
# #                 content_lines.append(content_match.group(1).strip())
# #             elif bullet_match:
# #                 # This line is bullet text
# #                 bullet_text_accumulator.append(bullet_match.group(1).strip())
# #             elif chart_match:
# #                 chart_text = chart_match.group(1).strip()
# #             elif table_match:
# #                 table_text = table_match.group(1).strip()
# #             else:
# #                 # If not recognized, treat as part of content
# #                 content_lines.append(line)
# #
# #         # Combine content lines into one string
# #         content_text = " ".join(content_lines).strip()
# #
# #         # Combine bullet lines into one big bullet string
# #         bullet_text_full = " ".join(bullet_text_accumulator).strip()
# #
# #         # Store in a dict
# #         slide_dict = {
# #             "slide_number": slide_num,
# #             "slide_subtitle": slide_subtitle,
# #             "content": content_text,
# #             "bullets": bullet_text_full,  # raw bullets, will split on '.' below
# #             "chart": chart_text,
# #             "table": table_text
# #         }
# #         slides_data.append(slide_dict)
# #
# #     # --- 5. Create slides in PPTX ---
# #     for sdata in slides_data:
# #         slide_layout = prs.slide_layouts[1]  # Title & Content layout
# #         slide = prs.slides.add_slide(slide_layout)
# #
# #         # Slide Title
# #         slide_title = f"Slide {sdata['slide_number']} ({sdata['slide_subtitle']})"
# #         slide.shapes.title.text = slide_title
# #
# #         # Content placeholder
# #         body_shape = slide.placeholders[1]
# #         text_frame = body_shape.text_frame
# #         text_frame.clear()
# #
# #         # 5a. Main Content
# #         if sdata["content"]:
# #             p = text_frame.add_paragraph()
# #             p.text = sdata["content"]
# #             p.level = 0  # non-bulleted paragraph
# #             p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
# #
# #         # 5b. Bullet Points
# #         if sdata["bullets"]:
# #             bullet_lines = split_bullets_on_period(sdata["bullets"])
# #             for bullet_line in bullet_lines:
# #                 p = text_frame.add_paragraph()
# #                 p.text = bullet_line
# #                 p.level = 1  # bullet
# #                 p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
# #
# #         # 5c. If there's a chart mention & we have a chart file
# #         if sdata["chart"] and "chart_filename" in data:
# #             chart_path = data["chart_filename"]
# #             if os.path.exists(chart_path):
# #                 left = Inches(1)
# #                 top = Inches(2)
# #                 width = Inches(6)
# #                 height = Inches(4)
# #                 slide.shapes.add_picture(chart_path, left, top, width, height)
# #             else:
# #                 print(f"Chart file not found: {chart_path}")
# #
# #         # 5d. If there's a table mention, you could place it or reference a table image similarly.
# #
# #     # --- 6. Save the PPT ---
# #     output_filename = "output2.pptx"
# #     prs.save(output_filename)
# #     print(f"Presentation saved as {output_filename}.")
# #
# #
# # # --------------------------------------------------------------------
# # # 6. MAIN FLOW
# # # --------------------------------------------------------------------
# #
# # if __name__ == "__main__":
# #     # 1) Ask user for file path
# #     file_path = input("Enter the file path: ")
# #
# #     # 2) Validate & convert file
# #     data = validate_and_convert(file_path)
# #     if "error" in data:
# #         print(f"Error processing file: {data['error']}")
# #         exit()
# #
# #     # 3) If data is tabular, get chart suggestion & generate chart
# #     if data.get("type") == "table":
# #         df = pd.DataFrame(data["content"])
# #         suggested_chart_type = gemini_recommend_chart_type(df)
# #         print(f"Gemini suggests a '{suggested_chart_type}' chart.")
# #
# #         chart_filename = "chart_0.png"
# #         plot_data_automatically(df, suggested_chart_type, output_name=chart_filename)
# #
# #         data["chart_filename"] = chart_filename
# #         data["chart_type"] = suggested_chart_type
# #
# #     # 4) Generate PPT
# #     generate_pptx_from_gemini(data, slides_title="My Gemini Presentation")
#
#
#
#
#
# import json
# import os
# import re
# import pandas as pd
# from bs4 import BeautifulSoup
# import seaborn as sns
# from PyPDF2 import PdfReader
# import google.generativeai as genai
# import matplotlib.pyplot as plt
# from processing_pdf import process_pdf  # Assuming you have a process_pdf function
# import pptx
# from pptx.util import Inches
# from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
#
# # --------------------------------------------------------------------
# # 0. CONFIGURE GEMINI
# # --------------------------------------------------------------------
# genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")  # Replace with your actual key
# model = genai.GenerativeModel("gemini-1.5-flash")
#
# # --------------------------------------------------------------------
# # 1. FILE PROCESSING FUNCTIONS
# # --------------------------------------------------------------------
#
# def process_text(text):
#     """Convert plain text into structured JSON."""
#     return {"type": "text", "content": text.strip()}
#
# def process_json(json_str):
#     """Parse JSON string and return a dictionary."""
#     try:
#         return json.loads(json_str)
#     except json.JSONDecodeError:
#         return {"error": "Invalid JSON format"}
#
# def process_table(file_path):
#     """Read tabular data (CSV/Excel) and convert it to JSON."""
#     try:
#         if file_path.endswith(".csv"):
#             df = pd.read_csv(file_path)
#         elif file_path.endswith(".xlsx"):
#             df = pd.read_excel(file_path)
#         else:
#             return {"error": "Unsupported table format"}
#         return {"type": "table", "content": df.to_dict(orient="records")}
#     except Exception as e:
#         return {"error": f"Failed to process table: {str(e)}"}
#
# def process_css(css_str):
#     """Parse CSS and return a JSON representation."""
#     try:
#         soup = BeautifulSoup(css_str, "html.parser")
#         rules = []
#         for rule in soup.text.split("}"):
#             if "{" in rule:
#                 selector, properties = rule.split("{", 1)
#                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
#         return {"type": "css", "content": rules}
#     except Exception as e:
#         return {"error": f"Failed to process CSS: {str(e)}"}
#
# def validate_and_convert(file_path):
#     """Validate and convert different file types to a dictionary representation."""
#     _, ext = os.path.splitext(file_path)
#     ext = ext.lower()
#
#     try:
#         if ext == ".txt":
#             with open(file_path, "r", encoding="utf-8") as f:
#                 return process_text(f.read())
#         elif ext == ".json":
#             with open(file_path, "r", encoding="utf-8") as f:
#                 return process_json(f.read())
#         elif ext in [".csv", ".xlsx"]:
#             return process_table(file_path)
#         elif ext == ".css":
#             with open(file_path, "r", encoding="utf-8") as f:
#                 return process_css(f.read())
#         elif ext == ".pdf":
#             pdf_result = process_pdf(file_path)  # Assuming you have a process_pdf function
#             if pdf_result.get("type") == "latex":
#                 return {"type": "latex", "content": pdf_result["content"]}
#             else:
#                 return pdf_result
#         else:
#             return {"error": "Unsupported file format"}
#     except Exception as e:
#         return {"error": f"Failed to process file: {str(e)}"}
#
# # --------------------------------------------------------------------
# # 2. CHART FUNCTIONS
# # --------------------------------------------------------------------
#
# def gemini_recommend_chart_type(table_df, purpose="visualize data"):
#     """
#     Calls Gemini to recommend a chart type for the given table.
#     """
#     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
#     sample_records = table_df.head(3).to_dict(orient="records")
#
#     prompt = f"""
# I have a table with columns: {columns_info}.
# Sample rows: {sample_records}.
# I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, pie, etc.)?
# Provide a short and concise answer.
# """
#     try:
#         response = model.generate_content(prompt)
#         return response.text.strip()
#     except Exception as e:
#         return f"Error calling Gemini API: {str(e)}"
#
# def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
#     """
#     Generates and saves a chart based on the table data and chart type suggestion.
#     """
#     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
#     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
#     suggestion_lower = chart_type_suggestion.lower()
#
#     plt.figure()
#     if suggestion_lower in ["line", "scatter"] and len(numeric_cols) >= 2:
#         if suggestion_lower == "line":
#             plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
#         else:
#             plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
#         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
#     elif suggestion_lower in ["hist", "distribution"] and len(numeric_cols) >= 1:
#         sns.histplot(table_df[numeric_cols[0]])
#         plt.title(f"{chart_type_suggestion.capitalize()} Plot")
#     elif suggestion_lower in ["bar", "column"] and len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
#         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
#         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
#     elif suggestion_lower in ["bar", "column"] and len(non_numeric_cols) >= 1:
#         sns.countplot(x=table_df[non_numeric_cols[0]])
#         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
#     else:
#         print(f"Unsupported or ambiguous chart type: {chart_type_suggestion}")
#         plt.close()
#         return
#
#     plt.tight_layout()
#     plt.savefig(output_name)
#     plt.close()
#
# # --------------------------------------------------------------------
# # 3. HELPER FOR SPLITTING BULLETS ON PERIOD
# # --------------------------------------------------------------------
#
# def split_bullets_on_period(bullet_text):
#     """
#     Splits bullet_text on a period + optional space => distinct bullet lines.
#     """
#     bullet_text = bullet_text.strip()
#     lines = re.split(r"\.\s+|\.$", bullet_text)
#     lines = [l.strip() for l in lines if l.strip()]
#     return lines
#
# # --------------------------------------------------------------------
# # 4. PPT GENERATION FUNCTION
# # --------------------------------------------------------------------
#
# def generate_pptx_from_gemini(data, slides_title="Presentation"):
#     """
#     1) Calls Gemini for a proposed slide structure.
#     2) Parses out each slide:
#         - The first line after "Slide X (Title):" is the actual slide heading.
#         - Additional lines are processed for content, bullet points, tables, etc.
#     3) If there's no table data, prints a message instead of generating a chart.
#     4) Saves the final PPTX.
#     """
#     prs = pptx.Presentation()
#
#     # --- Title Slide ---
#     title_slide_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(title_slide_layout)
#     slide.shapes.title.text = slides_title
#
#     # --- Build Gemini Prompt ---
#     gemini_prompt = f"""
# Create a PowerPoint presentation based on the following information:
#
# **Data:** {json.dumps(data, indent=4)}
#
# **Instructions:**
# - **Structure:** Suggest a logical slide order and content for each slide.
# - **Visualizations:** Mention any chart or table if needed ("- Chart:" or "- Table: ...").
# - **Text Formatting:** Provide headings, bullet points, and concise text for each slide.
#
# Use this format:
# - **Slide 1 (Title): Actual Slide Title**
#   - Content: ...
#   - Bullet Points: ...
#   - Table: ...
#   - Chart: ...
# - **Slide 2 (Title): Some Other Title**
#   ...
#     """
#     try:
#         response = model.generate_content(gemini_prompt)
#         gemini_text = response.text.strip()
#         print("[Gemini Slide Structure Suggestion]:")
#         print(gemini_text)
#     except Exception as e:
#         print(f"Error calling Gemini API: {str(e)}")
#         return
#
#     # --- Regex to find each slide ---
#     slide_pattern = (
#         r"- \*\*Slide\s+(\d+)\s*\(([^)]+)\):"   # e.g., "- **Slide 1 (Title):"
#         r"(.*?)"                                # capture the rest (non-greedy)
#         r"(?=- \*\*Slide|\Z)"                   # until next slide or end
#     )
#     matches = re.findall(slide_pattern, gemini_text, flags=re.DOTALL)
#
#     slides_data = []
#     for slide_num, slide_subtitle, slide_body in matches:
#         slide_num = slide_num.strip()
#         slide_subtitle = slide_subtitle.strip()
#         slide_body = slide_body.strip()
#
#         # Now, we want the *first line* of slide_body to be the actual slide title
#         # e.g. "SmartCamID: Real-Time Camera Stream Encoding..." if the line is there
#         lines = [ln.strip() for ln in slide_body.splitlines() if ln.strip()]
#
#         if not lines:
#             # If there's no extra lines, fallback to using slide_subtitle
#             actual_slide_title = slide_subtitle
#             body_lines = []
#         else:
#             # The first line is the real slide heading
#             actual_slide_title = lines[0]
#             body_lines = lines[1:]  # the rest is body
#
#         # We'll parse the "body_lines" for - Content:, - Bullet Points:, etc.
#         content_lines = []
#         bullet_text_accumulator = []
#         chart_text = ""
#         table_text = ""
#
#         for line in body_lines:
#             c_match = re.match(r"^- Content:\s*(.*)$", line, re.IGNORECASE)
#             b_match = re.match(r"^- Bullet Points?:\s*(.*)$", line, re.IGNORECASE)
#             chart_match = re.match(r"^- Chart:\s*(.*)$", line, re.IGNORECASE)
#             table_match = re.match(r"^- Table:\s*(.*)$", line, re.IGNORECASE)
#
#             if c_match:
#                 content_lines.append(c_match.group(1).strip())
#             elif b_match:
#                 bullet_text_accumulator.append(b_match.group(1).strip())
#             elif chart_match:
#                 chart_text = chart_match.group(1).strip()
#             elif table_match:
#                 table_text = table_match.group(1).strip()
#             else:
#                 # If it doesn't match, treat it as extra content
#                 content_lines.append(line)
#
#         final_content = " ".join(content_lines).strip()
#         final_bullets = " ".join(bullet_text_accumulator).strip()
#
#         slides_data.append({
#             "slide_number": slide_num,
#             "slide_subtitle": slide_subtitle,
#             "slide_title": actual_slide_title,
#             "content": final_content,
#             "bullets": final_bullets,
#             "chart": chart_text,
#             "table": table_text
#         })
#
#     # --- Create each slide ---
#     for sdata in slides_data:
#         # We'll put the user-chosen 'slide_title' as the actual PPT slide title
#         slide_layout = prs.slide_layouts[1]  # Title & Content
#         slide = prs.slides.add_slide(slide_layout)
#
#         # Put the real title
#         slide.shapes.title.text = sdata["slide_title"]
#
#         # Content placeholder
#         body_shape = slide.placeholders[1]
#         text_frame = body_shape.text_frame
#         text_frame.clear()
#
#         # Paragraph for content
#         if sdata["content"]:
#             p = text_frame.add_paragraph()
#             p.text = sdata["content"]
#             p.level = 0
#
#         # Paragraphs for bullet points
#         if sdata["bullets"]:
#             bullet_lines = split_bullets_on_period(sdata["bullets"])
#             for b in bullet_lines:
#                 p = text_frame.add_paragraph()
#                 p.text = b
#                 p.level = 1
#
#         # If there's a chart mention & we actually generated one
#         if sdata["chart"] and "chart_filename" in data:
#             chart_path = data["chart_filename"]
#             if os.path.exists(chart_path):
#                 left = Inches(1)
#                 top = Inches(2)
#                 width = Inches(6)
#                 height = Inches(4)
#                 slide.shapes.add_picture(chart_path, left, top, width, height)
#             else:
#                 print("Chart file not found:", chart_path)
#
#         # If there's a table mention, handle similarly as needed
#
#     # --- Save the PPTX ---
#     output_filename = "output2.pptx"
#     prs.save(output_filename)
#     print(f"Presentation saved as {output_filename}.")
#
# # --------------------------------------------------------------------
# # 5. MAIN
# # --------------------------------------------------------------------
#
# if __name__ == "__main__":
#     file_path = input("Enter the file path: ")
#     data = validate_and_convert(file_path)
#
#     if "error" in data:
#         print("Error processing file:", data["error"])
#         exit()
#
#     # If there's no table data => "No table data found."
#     # Otherwise, generate a chart.
#     if data.get("type") == "table":
#         df = pd.DataFrame(data["content"])
#         if df.empty:
#             print("No table data found.")
#         else:
#             chart_type = gemini_recommend_chart_type(df)
#             print("Gemini suggests chart type:", chart_type)
#             plot_data_automatically(df, chart_type, "chart_0.png")
#             data["chart_filename"] = "chart_0.png"
#     else:
#         print("No table data found.")
#
#     # Finally, generate the PPT using Gemini
#     generate_pptx_from_gemini(data, slides_title="My Gemini Presentation")



# import json
# import os
# import re
# import pandas as pd
# from bs4 import BeautifulSoup
# import seaborn as sns
# from PyPDF2 import PdfReader
# import google.generativeai as genai
# import matplotlib.pyplot as plt
# from processing_pdf import process_pdf  # Assuming you have a process_pdf function
# import pptx
# from pptx.util import Inches
# from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
#
# # --------------------------------------------------------------------
# # 0. CONFIGURE GEMINI
# # --------------------------------------------------------------------
# genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")  # Replace with your actual key
# model = genai.GenerativeModel("gemini-1.5-flash")
#
# # --------------------------------------------------------------------
# # 1. FILE PROCESSING FUNCTIONS
# # --------------------------------------------------------------------
#
# def process_text(text):
#     """Convert plain text into structured JSON."""
#     return {"type": "text", "content": text.strip()}
#
# def process_json(json_str):
#     """Parse JSON string and return a dictionary."""
#     try:
#         return json.loads(json_str)
#     except json.JSONDecodeError:
#         return {"error": "Invalid JSON format"}
#
# def process_table(file_path):
#     """Read tabular data (CSV/Excel) and convert it to JSON."""
#     try:
#         if file_path.endswith(".csv"):
#             df = pd.read_csv(file_path)
#         elif file_path.endswith(".xlsx"):
#             df = pd.read_excel(file_path)
#         else:
#             return {"error": "Unsupported table format"}
#         return {"type": "table", "content": df.to_dict(orient="records")}
#     except Exception as e:
#         return {"error": f"Failed to process table: {str(e)}"}
#
# def process_css(css_str):
#     """Parse CSS and return a JSON representation."""
#     try:
#         soup = BeautifulSoup(css_str, "html.parser")
#         rules = []
#         for rule in soup.text.split("}"):
#             if "{" in rule:
#                 selector, properties = rule.split("{", 1)
#                 rules.append({"selector": selector.strip(), "properties": properties.strip()})
#         return {"type": "css", "content": rules}
#     except Exception as e:
#         return {"error": f"Failed to process CSS: {str(e)}"}
#
# def validate_and_convert(file_path):
#     """Validate and convert different file types to a dictionary representation."""
#     _, ext = os.path.splitext(file_path)
#     ext = ext.lower()
#
#     try:
#         if ext == ".txt":
#             with open(file_path, "r", encoding="utf-8") as f:
#                 return process_text(f.read())
#         elif ext == ".json":
#             with open(file_path, "r", encoding="utf-8") as f:
#                 return process_json(f.read())
#         elif ext in [".csv", ".xlsx"]:
#             return process_table(file_path)
#         elif ext == ".css":
#             with open(file_path, "r", encoding="utf-8") as f:
#                 return process_css(f.read())
#         elif ext == ".pdf":
#             pdf_result = process_pdf(file_path)  # Assuming you have a process_pdf function
#             if pdf_result.get("type") == "latex":
#                 return {"type": "latex", "content": pdf_result["content"]}
#             else:
#                 return pdf_result
#         else:
#             return {"error": "Unsupported file format"}
#     except Exception as e:
#         return {"error": f"Failed to process file: {str(e)}"}
#
# # --------------------------------------------------------------------
# # 2. CHART FUNCTIONS
# # --------------------------------------------------------------------
#
# def gemini_recommend_chart_type(table_df, purpose="visualize data"):
#     """
#     Calls Gemini to recommend a chart type for the given table.
#     """
#     columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
#     sample_records = table_df.head(3).to_dict(orient="records")
#
#     prompt = f"""
# I have a table with columns: {columns_info}.
# Sample rows: {sample_records}.
# I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, pie etc.)?
# Provide a short and concise answer.
# """
#     try:
#         response = model.generate_content(prompt)
#         return response.text.strip()
#     except Exception as e:
#         return f"Error calling Gemini API: {str(e)}"
#
# def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
#     """
#     Generates and saves a chart based on the table data and chart type suggestion.
#     """
#     numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
#     non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
#     suggestion_lower = chart_type_suggestion.lower()
#
#     plt.figure()
#     if suggestion_lower in ["line", "scatter"] and len(numeric_cols) >= 2:
#         if suggestion_lower == "line":
#             plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
#         else:
#             plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
#         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
#     elif suggestion_lower in ["hist", "distribution"] and len(numeric_cols) >= 1:
#         sns.histplot(table_df[numeric_cols[0]])
#         plt.title(f"{chart_type_suggestion.capitalize()} Plot")
#     elif suggestion_lower in ["bar", "column"] and len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
#         plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
#         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
#     elif suggestion_lower in ["bar", "column"] and len(non_numeric_cols) >= 1:
#         sns.countplot(x=table_df[non_numeric_cols[0]])
#         plt.title(f"{chart_type_suggestion.capitalize()} Chart")
#     else:
#         print(f"Unsupported or ambiguous chart type: {chart_type_suggestion}")
#         plt.close()
#         return
#
#     plt.tight_layout()
#     plt.savefig(output_name)
#     plt.close()
#
# # --------------------------------------------------------------------
# # 3. HELPER FOR SPLITTING BULLETS ON PERIOD
# # --------------------------------------------------------------------
#
# def split_bullets_on_period(bullet_text):
#     """
#     Splits bullet_text on a period + optional space => distinct bullet lines.
#     """
#     bullet_text = bullet_text.strip()
#     lines = re.split(r"\.\s+|\.$", bullet_text)
#     lines = [l.strip() for l in lines if l.strip()]
#     return lines
#
# # --------------------------------------------------------------------
# # 4. HELPERS FOR REFERENCES SLIDE CLEANUP
# # --------------------------------------------------------------------
#
# def cleanup_unwanted_references_text(text):
#     """
#     Removes placeholder or extra lines from the references slide content.
#     For example, remove lines like:
#       "Remember to replace placeholder text with actual data..."
#     Adjust the patterns as needed to remove more unwanted lines.
#     """
#     # Remove any mention of "Remember to replace placeholder text..." case-insensitively
#     text = re.sub(r"(?i)remember to replace placeholder.*", "", text).strip()
#     # Remove mention of "The LaTeX document provides ..." if present
#     text = re.sub(r"(?i)the latex document provides.*", "", text).strip()
#     return text
#
# # --------------------------------------------------------------------
# # 5. PPT GENERATION FUNCTION
# # --------------------------------------------------------------------
#
# def generate_pptx_from_gemini(data, slides_title="Presentation"):
#     """
#     1) Calls Gemini for a proposed slide structure.
#     2) Parses out each slide:
#         - The first line after "Slide X (Title):" is the real slide heading.
#         - Additional lines are processed for content, bullet points, tables, charts.
#     3) Removes placeholder references text from the "References" slide.
#     4) If there's no table data, prints a message instead of generating a chart.
#     5) Saves the final PPTX.
#     """
#     prs = pptx.Presentation()
#
#     # --- 1. Title Slide ---
#     title_slide_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(title_slide_layout)
#     slide.shapes.title.text = slides_title
#
#     # --- 2. Build Gemini Prompt ---
#     gemini_prompt = f"""
# Create a PowerPoint presentation based on the following information:
#
# **Data:** {json.dumps(data, indent=4)}
#
# **Instructions:**
# - **Structure:** Provide a clear and concise structure, specifying the content for each slide. Use separate prompts for content with bullet points and content for paragraphs.
# - **Slide Titles:** Use colons (:) to separate slide numbers and titles (e.g., "Slide 1: Introduction").
# - **Visualizations:** Mention any chart or table if needed ("- Chart:" or "- Table: ...").
# - **Text Formatting:** Specify bullet points or paragraphs for content as needed.
#
# **Example Prompts:**
#   - Slide 1: Introduction
#     - Content (no bullet points): This is the opening paragraph of the presentation.
#
#   - Slide 2: Key Findings
#     - Bullet points:
#       - Finding 1
#       - Finding 2
#
#   - Slide 3: Conclusion
#     - Content (no bullet points): Briefly summarize the key takeaways.
#
# **Note:** Break down information into multiple slides if it improves clarity.
#   ...
#     """
#     try:
#         response = model.generate_content(gemini_prompt)
#         gemini_text = response.text.strip()
#         print("[Gemini Slide Structure Suggestion]:")
#         print(gemini_text)
#     except Exception as e:
#         print(f"Error calling Gemini API: {str(e)}")
#         return
#
#     # --- 3. Parse each slide using regex ---
#     slide_pattern = (
#         r"- \*\*Slide\s+(\d+)\s*\(([^)]+)\):"  # e.g., "- **Slide 1 (Title):"
#         r"(.*?)"                               # capture the rest
#         r"(?=- \*\*Slide|\Z)"                  # until next slide or end
#     )
#     matches = re.findall(slide_pattern, gemini_text, flags=re.DOTALL)
#
#     slides_data = []
#     for slide_num, slide_subtitle, slide_body in matches:
#         slide_num = slide_num.strip()
#         slide_subtitle = slide_subtitle.strip()
#         slide_body = slide_body.strip()
#
#         # Split the body into lines, ignoring blank lines
#         lines = [ln.strip() for ln in slide_body.splitlines() if ln.strip()]
#
#         # The first line is the real PPT slide title
#         if lines:
#             actual_slide_title = lines[0]
#             body_lines = lines[1:]
#         else:
#             actual_slide_title = slide_subtitle
#             body_lines = []
#
#         content_lines = []
#         bullet_text_accumulator = []
#         chart_text = ""
#         table_text = ""
#
#         for line in body_lines:
#             c_match = re.match(r"^- Content:\s*(.*)$", line, re.IGNORECASE)
#             b_match = re.match(r"^- Bullet Points?:\s*(.*)$", line, re.IGNORECASE)
#             chart_match = re.match(r"^- Chart:\s*(.*)$", line, re.IGNORECASE)
#             table_match = re.match(r"^- Table:\s*(.*)$", line, re.IGNORECASE)
#
#             if c_match:
#                 content_lines.append(c_match.group(1).strip())
#             elif b_match:
#                 bullet_text_accumulator.append(b_match.group(1).strip())
#             elif chart_match:
#                 chart_text = chart_match.group(1).strip()
#             elif table_match:
#                 table_text = table_match.group(1).strip()
#             else:
#                 # If not a recognized marker, treat it as additional content
#                 content_lines.append(line)
#
#         # Combine lines into final content
#         final_content = " ".join(content_lines).strip()
#         final_bullets = " ".join(bullet_text_accumulator).strip()
#
#         # If the title or subtitle indicates "References", do cleanup
#         if re.search(r"(?i)references", actual_slide_title):
#             # Remove placeholders from final_content
#             final_content = cleanup_unwanted_references_text(final_content)
#
#         slides_data.append({
#             "slide_number": slide_num,
#             "slide_subtitle": slide_subtitle,
#             "slide_title": actual_slide_title,
#             "content": final_content,
#             "bullets": final_bullets,
#             "chart": chart_text,
#             "table": table_text
#         })
#
#     # --- 4. Create slides in PPTX ---
#     for sdata in slides_data:
#         slide_layout = prs.slide_layouts[1]  # Title & Content
#         slide = prs.slides.add_slide(slide_layout)
#
#         # Use the "real" title in slide_title
#         slide.shapes.title.text = sdata["slide_title"]
#
#         # Content placeholder
#         body_shape = slide.placeholders[1]
#         text_frame = body_shape.text_frame
#         text_frame.clear()
#
#         # Content paragraph
#         if sdata["content"]:
#             p = text_frame.add_paragraph()
#             p.text = sdata["content"]
#             p.level = 0
#
#         # Bullet points
#         if sdata["bullets"]:
#             bullet_lines = split_bullets_on_period(sdata["bullets"])
#             for b_line in bullet_lines:
#                 p = text_frame.add_paragraph()
#                 p.text = b_line
#                 p.level = 1
#
#         # If there's a chart mention & we have a chart file
#         if sdata["chart"] and "chart_filename" in data:
#             chart_path = data["chart_filename"]
#             if os.path.exists(chart_path):
#                 left = Inches(1)
#                 top = Inches(2)
#                 width = Inches(6)
#                 height = Inches(4)
#                 slide.shapes.add_picture(chart_path, left, top, width, height)
#             else:
#                 print("Chart file not found:", chart_path)
#
#         # If there's a table mention, handle similarly, or skip for now
#
#     # --- 5. Save the PPTX ---
#     output_filename = "output2.pptx"
#     prs.save(output_filename)
#     print(f"Presentation saved as {output_filename}.")
#
# # --------------------------------------------------------------------
# # 6. MAIN FLOW
# # --------------------------------------------------------------------
#
# if __name__ == "__main__":
#     file_path = input("Enter the file path: ")
#     data = validate_and_convert(file_path)
#
#     if "error" in data:
#         print("Error processing file:", data["error"])
#         exit()
#
#     # If the data is a table, generate a chart if possible
#     if data.get("type") == "table":
#         df = pd.DataFrame(data["content"])
#         if df.empty:
#             print("No table data found.")
#         else:
#             chart_type = gemini_recommend_chart_type(df)
#             print("Gemini suggests chart type:", chart_type)
#             plot_data_automatically(df, chart_type, "chart_0.png")
#             data["chart_filename"] = "chart_0.png"
#     else:
#         print("No table data found.")
#
#     # Generate the PPT
#     generate_pptx_from_gemini(data, slides_title="My Gemini Presentation")


import json
import os
import re
import pandas as pd
from bs4 import BeautifulSoup
import seaborn as sns
from PyPDF2 import PdfReader
import google.generativeai as genai
import matplotlib.pyplot as plt
from processing_pdf import process_pdf  # Assuming you have a process_pdf function
import pptx
from pptx.util import Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

# --------------------------------------------------------------------
# 0. CONFIGURE GEMINI
# --------------------------------------------------------------------
genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")  # Replace with your actual key
model = genai.GenerativeModel("gemini-1.5-flash")

# --------------------------------------------------------------------
# 1. FILE PROCESSING FUNCTIONS
# --------------------------------------------------------------------

def process_text(text):
    """Convert plain text into structured JSON."""
    return {"type": "text", "content": text.strip()}

def process_json(json_str):
    """Parse JSON string and return a dictionary."""
    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        return {"error": "Invalid JSON format"}

def process_table(file_path):
    """Read tabular data (CSV/Excel) and convert it to JSON."""
    try:
        if file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
        elif file_path.endswith(".xlsx"):
            df = pd.read_excel(file_path)
        else:
            return {"error": "Unsupported table format"}
        return {"type": "table", "content": df.to_dict(orient="records")}
    except Exception as e:
        return {"error": f"Failed to process table: {str(e)}"}

def process_css(css_str):
    """Parse CSS and return a JSON representation."""
    try:
        soup = BeautifulSoup(css_str, "html.parser")
        rules = []
        for rule in soup.text.split("}"):
            if "{" in rule:
                selector, properties = rule.split("{", 1)
                rules.append({"selector": selector.strip(), "properties": properties.strip()})
        return {"type": "css", "content": rules}
    except Exception as e:
        return {"error": f"Failed to process CSS: {str(e)}"}

def validate_and_convert(file_path):
    """Validate and convert different file types to a dictionary representation."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return process_text(f.read())
        elif ext == ".json":
            with open(file_path, "r", encoding="utf-8") as f:
                return process_json(f.read())
        elif ext in [".csv", ".xlsx"]:
            return process_table(file_path)
        elif ext == ".css":
            with open(file_path, "r", encoding="utf-8") as f:
                return process_css(f.read())
        elif ext == ".pdf":
            pdf_result = process_pdf(file_path)  # Assuming you have a process_pdf function
            if pdf_result.get("type") == "latex":
                return {"type": "latex", "content": pdf_result["content"]}
            else:
                return pdf_result
        else:
            return {"error": "Unsupported file format"}
    except Exception as e:
        return {"error": f"Failed to process file: {str(e)}"}

# --------------------------------------------------------------------
# 2. CHART FUNCTIONS
# --------------------------------------------------------------------

def gemini_recommend_chart_type(table_df, purpose="visualize data"):
    """
    Calls Gemini to recommend a chart type for the given table.
    """
    columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
    sample_records = table_df.head(3).to_dict(orient="records")

    prompt = f"""
I have a table with columns: {columns_info}.
Sample rows: {sample_records}.
I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, pie etc.)?
Provide a short and concise answer.
"""
    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        return f"Error calling Gemini API: {str(e)}"

def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
    """
    Generates and saves a chart based on the table data and chart type suggestion.
    """
    numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
    non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
    suggestion_lower = chart_type_suggestion.lower()

    plt.figure()
    if suggestion_lower in ["line", "scatter"] and len(numeric_cols) >= 2:
        if suggestion_lower == "line":
            plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
        else:
            plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
        plt.title(f"{chart_type_suggestion.capitalize()} Chart")
    elif suggestion_lower in ["hist", "distribution"] and len(numeric_cols) >= 1:
        sns.histplot(table_df[numeric_cols[0]])
        plt.title(f"{chart_type_suggestion.capitalize()} Plot")
    elif suggestion_lower in ["bar", "column"] and len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
        plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
        plt.title(f"{chart_type_suggestion.capitalize()} Chart")
    elif suggestion_lower in ["bar", "column"] and len(non_numeric_cols) >= 1:
        sns.countplot(x=table_df[non_numeric_cols[0]])
        plt.title(f"{chart_type_suggestion.capitalize()} Chart")
    else:
        print(f"Unsupported or ambiguous chart type: {chart_type_suggestion}")
        plt.close()
        return

    plt.tight_layout()
    plt.savefig(output_name)
    plt.close()

# --------------------------------------------------------------------
# 3. HELPER FOR SPLITTING BULLETS ON PERIOD
# --------------------------------------------------------------------

def split_bullets_on_period(bullet_text):
    """
    Splits bullet_text on a period + optional space => distinct bullet lines.
    """
    bullet_text = bullet_text.strip()
    lines = re.split(r"\.\s+|\.$", bullet_text)
    lines = [l.strip() for l in lines if l.strip()]
    return lines

# --------------------------------------------------------------------
# 4. HELPERS FOR REFERENCES SLIDE CLEANUP
# --------------------------------------------------------------------

def cleanup_unwanted_references_text(text):
    """
    Removes placeholder or extra lines from the references slide content.
    For example, remove lines like:
      "Remember to replace placeholder text with actual data..."
    Adjust the patterns as needed to remove more unwanted lines.
    """
    # Remove any mention of "Remember to replace placeholder text..." case-insensitively
    text = re.sub(r"(?i)remember to replace placeholder.*", "", text).strip()
    # Remove mention of "The LaTeX document provides ..." if present
    text = re.sub(r"(?i)the latex document provides.*", "", text).strip()
    return text

# --------------------------------------------------------------------
# 5. PPT GENERATION FUNCTION (UPDATED)
# --------------------------------------------------------------------
def generate_pptx_from_gemini(data, slides_title="Presentation"):
    """
    1) Calls Gemini for a proposed slide structure in the new instructions style:
       Slide 1: Introduction
       Content (no bullet points): ...
       Bullet points:
       - ...
       - ...
    2) Parses by searching for "Slide X: Title" blocks, then reading lines for content,
       bullet points, and optional Chart/Table mentions.
    3) Removes '**' from any text before placing it in PPT.
    4) Removes placeholder references if it's a References slide.
    5) Saves the final PPTX.
    """
    prs = pptx.Presentation()
    print("here ----------------------------------------------------")
    print(data)
    # --- Create Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = slides_title

    # --- 1. Build Gemini Prompt ---
    gemini_prompt = f"""
Create a PowerPoint presentation based on the following information:

**Data:** {json.dumps(data, indent=4)}

**Instructions:**
- **Structure:** Provide a clear and concise structure, specifying the content for each slide. Dont give unecessary "*"
- **Slide Titles:** Use colons (:) to separate slide numbers and titles (e.g., "Slide 1: Introduction").
- **Text Formatting:** 
   - If you have a paragraph of text, write: 
       Content (no bullet points): <text here>
   - If you have bullet points, write:
       Bullet points:
       - ...
       - ...
   -Very impornt - do not tell me to take text from anywhere just give what will be going in each slide i dont
    want to copy anything from the pdf u give what all is to be added here itself         
- **Visualizations:** If needed, mention them with lines like "- Chart:" or "- Table: ...".
- **Example**:
  Slide 1: Introduction
  Content (no bullet points): This is the opening paragraph.

  Slide 2: Key Findings
  Bullet points:
  - Finding 1
  - Finding 2

Break down information into multiple slides if it improves clarity.
"""

    try:
        response = model.generate_content(gemini_prompt)
        gemini_text = response.text.strip()
        print("[Gemini Slide Structure Suggestion]:")
        print(gemini_text)
    except Exception as e:
        print(f"Error calling Gemini API: {str(e)}")
        return

    # --- 2. Parse for "Slide X: Title" blocks ---
    slide_block_pattern = r"(Slide\s+(\d+)\s*:\s*(.*?))(?=Slide\s+\d+|$)"
    raw_blocks = re.findall(slide_block_pattern, gemini_text, flags=re.DOTALL)

    slides_data = []
    for full_block, slide_num_str, after_colon_title in raw_blocks:
        # Convert slide_num_str to int if possible
        try:
            slide_num = int(slide_num_str.strip())
        except ValueError:
            slide_num = 999  # fallback if parsing fails

        # The remainder after "Slide X: " is the nominal slide title
        lines = full_block.splitlines()
        # lines[0] is "Slide X: Title"
        block_lines = lines[1:]  # everything after that line

        # Slide title = the part after the colon, up to newline
        slide_title = after_colon_title.split("\n")[0].strip()

        # We'll parse subsequent lines for:
        #   Content (no bullet points): ...
        #   Bullet points:
        #       - item
        #       - item
        #   - Chart:
        #   - Table:
        content_str = ""
        bullet_points = []
        chart_info = ""
        table_info = ""

        # We'll track a 'bullet mode' once we see "Bullet points:"
        in_bullet_mode = False

        for line in block_lines:
            txt = line.strip()
            if not txt:
                # empty line => end bullet mode
                in_bullet_mode = False
                continue

            # Check patterns
            content_match = re.match(
                r"^Content\s*\(no\s*bullet\s*points\)\s*:\s*(.*)$", txt, re.IGNORECASE
            )
            bulletpts_match = re.match(r"^Bullet\s*points\s*:\s*$", txt, re.IGNORECASE)
            chart_match = re.match(r"^-?\s*Chart:\s*(.*)$", txt, re.IGNORECASE)
            table_match = re.match(r"^-?\s*Table:\s*(.*)$", txt, re.IGNORECASE)

            if content_match:
                content_str += content_match.group(1).strip() + " "
                in_bullet_mode = False
            elif bulletpts_match:
                in_bullet_mode = True
            elif chart_match:
                chart_info = chart_match.group(1).strip()
                in_bullet_mode = False
            elif table_match:
                table_info = table_match.group(1).strip()
                in_bullet_mode = False
            else:
                # Could be a bullet item if in bullet mode and line starts with "- "
                if in_bullet_mode and txt.startswith("-"):
                    bullet_points.append(txt[1:].strip())
                else:
                    # Otherwise treat as content
                    content_str += txt + " "

        final_content = content_str.strip()

        # If the slide title is "References", clean placeholders
        if re.search(r"(?i)references", slide_title):
            final_content = cleanup_unwanted_references_text(final_content)

        # ---------- Remove "**" from everything ----------
        slide_title = slide_title.replace("**", "")
        final_content = final_content.replace("**", "")
        bullet_points = [bp.replace("**", "") for bp in bullet_points]
        chart_info = chart_info.replace("**", "")
        table_info = table_info.replace("**", "")
        # -------------------------------------------------

        slides_data.append({
            "slide_num": slide_num,
            "slide_title": slide_title,
            "content": final_content,
            "bullets": bullet_points,
            "chart": chart_info,
            "table": table_info
        })

    # Sort slides_data by slide_num
    slides_data.sort(key=lambda x: x["slide_num"])

    # --- 3. Create slides in PPTX ---
    for sdata in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title & Content
        slide = prs.slides.add_slide(slide_layout)

        # Slide Title
        slide.shapes.title.text = sdata["slide_title"]

        # Content placeholder
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        # Add content paragraph
        if sdata["content"]:
            p = text_frame.add_paragraph()
            p.text = sdata["content"]
            p.level = 0

        # Add bullet points
        for bullet in sdata["bullets"]:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 1

        # If there's a chart mention and we have a chart file
        if sdata["chart"] and "chart_filename" in data:
            chart_path = data["chart_filename"]
            if os.path.exists(chart_path):
                left = Inches(1)
                top = Inches(2)
                width = Inches(6)
                height = Inches(4)
                slide.shapes.add_picture(chart_path, left, top, width, height)
            else:
                print("Chart file not found:", chart_path)

    # --- 4. Save the PPTX ---
    output_filename = "output2.pptx"
    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}.")


# --------------------------------------------------------------------
# 6. MAIN FLOW
# --------------------------------------------------------------------

if __name__ == "__main__":
    file_path = input("Enter the file path: ")
    data = validate_and_convert(file_path)

    if "error" in data:
        print("Error processing file:", data["error"])
        exit()

    # If the data is a table, generate a chart if possible
    if data.get("type") == "table":
        df = pd.DataFrame(data["content"])
        if df.empty:
            print("No table data found.")
        else:
            chart_type = gemini_recommend_chart_type(df)
            print("Gemini suggests chart type:", chart_type)
            plot_data_automatically(df, chart_type, "chart_0.png")
            data["chart_filename"] = "chart_0.png"
    else:
        print("No table data found.")

    # Generate the PPT
    generate_pptx_from_gemini(data, slides_title="My Gemini Presentation")


# --------------------------------------------------------------------
# 6. MAIN FLOW
# --------------------------------------------------------------------

if __name__ == "__main__":
    file_path = input("Enter the file path: ")
    data = validate_and_convert(file_path)

    if "error" in data:
        print("Error processing file:", data["error"])
        exit()

    # If the data is a table, generate a chart if possible
    if data.get("type") == "table":
        df = pd.DataFrame(data["content"])
        if df.empty:
            print("No table data found.")
        else:
            chart_type = gemini_recommend_chart_type(df)
            print("Gemini suggests chart type:", chart_type)
            plot_data_automatically(df, chart_type, "chart_0.png")
            data["chart_filename"] = "chart_0.png"
    else:
        print("No table data found.")

    # Generate the PPT
    generate_pptx_from_gemini(data, slides_title="My Gemini Presentation")
