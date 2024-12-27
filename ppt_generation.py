import json
import os
import re
import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

import plot_code
from gemini_config import model
from helpers import cleanup_unwanted_references_text
import PIL.Image

def generate_pptx_from_gemini(data, slides_title="Presentation"):
    """
    1) Calls Gemini for a proposed slide structure.
    2) Parses each "Slide X: Title" block for content, bullet points,
       and a 'Chart:' line that (ideally) has something like '(filename.png)'.
    3) Uses a blank layout and manually places:
       - A top-centered title
       - All text in the upper half (or so) of the slide
       - The chart at the bottom if the slide mentions a chart
    4) Saves the final PPTX as output2.pptx.
    """

    prs = pptx.Presentation()

    print("Generating charts...")
    # ----------------------------------------------------------------
    # 1) Generate or retrieve the list of chart files from plot_code
    #    e.g. ["temperature.png", "rainfall.png", "sales_vs_budget.png", "sales.png"]
    # ----------------------------------------------------------------
    chart_file_list = plot_code.main()
    data["chart_filename"] = chart_file_list

    # print("Data after charts:\n", data)

    # ----------------------------------------------------------------
    # 2) Build the Gemini prompt (no extra "title slide" from us)
    # ----------------------------------------------------------------
    doc_summary = data.get("document_summary", "")
    gemini_prompt = f"""
    Create a PowerPoint presentation based on the following information:

    **Data:** 
    {json.dumps(data, indent=4)}

    **Document Summary (if any):** 
    {doc_summary}

    **Instructions:**
    - **Structure:** Provide a clear, concise slide structure.
    - **Slide Titles:** Use "Slide X: Title".
    - **Text Formatting:**
      - For a paragraph: "Content (no bullet points): <text>"
      - For bullet points:
          Bullet points:
          - ...
          - ...
      - If referencing the chart, mention:
          - Chart: Some text (my_chart.png)
        (Put the chart filename in parentheses so we can parse it.)
    - Summarize the text and mention key points.
    - ***IMPORTANT***: Provide EXACT text for each slide. 
      Do not just say "copy from data".
    - ***IMPORTANT***: We have chart files: {data.get("chart_filename", "No chart")} 
      - mention how they fit in the presentation.

    Break down the information into multiple slides if it improves clarity.
    """

    # ----------------------------------------------------------------
    # 3) Call Gemini to get the slide structure
    # ----------------------------------------------------------------
    try:
        response = model.generate_content(gemini_prompt)
        gemini_text = response.text.strip()
        print("[Gemini Slide Structure Suggestion]:")
        print(gemini_text)
    except Exception as e:
        print(f"Error calling Gemini API: {str(e)}")
        return

    # ----------------------------------------------------------------
    # 4) Parse the Gemini text for slides: "Slide X: Title"
    # ----------------------------------------------------------------
    slide_block_pattern = r"(Slide\s+(\d+)\s*:\s*(.*?))(?=Slide\s+\d+|$)"
    raw_blocks = re.findall(slide_block_pattern, gemini_text, flags=re.DOTALL)

    slides_data = []
    for full_block, slide_num_str, after_colon_title in raw_blocks:
        try:
            slide_num = int(slide_num_str.strip())
        except ValueError:
            slide_num = 999

        lines = full_block.splitlines()
        block_lines = lines[1:]  # everything after "Slide X: Title"
        slide_title = after_colon_title.split("\n")[0].strip()

        content_str = ""
        bullet_points = []
        chart_info = ""
        table_info = ""
        in_bullet_mode = False

        for line in block_lines:
            txt = line.strip()
            if not txt:
                in_bullet_mode = False
                continue

            content_match = re.match(
                r"^Content\s*\(no\s*bullet\s*points\)\s*:\s*(.*)$", txt, re.IGNORECASE
            )
            bulletpts_match = re.match(r"^Bullet\s*points\s*:\s*$", txt, re.IGNORECASE)
            chart_match = re.match(r"^-?\s*Chart:\s*(.*)$", txt, re.IGNORECASE)
            table_match = re.match(r"^-?\s*Table:\s*(.*)$", txt, re.IGNORECASE)

            if content_match:
                content_str += content_match.group(1).strip() + " "
                in_bullet_mode = False
            elif bulletpts_match:
                in_bullet_mode = True
            elif chart_match:
                chart_info = chart_match.group(1).strip()
                in_bullet_mode = False
            elif table_match:
                table_info = table_match.group(1).strip()
                in_bullet_mode = False
            else:
                if in_bullet_mode and txt.startswith("-"):
                    bullet_points.append(txt[1:].strip())
                else:
                    content_str += txt + " "

        final_content = content_str.strip()

        # If the slide title is "References", remove placeholders
        if re.search(r"(?i)references", slide_title):
            final_content = cleanup_unwanted_references_text(final_content)

        # Remove "**" from text
        slide_title = slide_title.replace("**", "")
        final_content = final_content.replace("**", "")
        bullet_points = [bp.replace("**", "") for bp in bullet_points]
        chart_info = chart_info.replace("**", "")
        table_info = table_info.replace("**", "")

        slides_data.append({
            "slide_num": slide_num,
            "slide_title": slide_title,
            "content": final_content,
            "bullets": bullet_points,
            "chart": chart_info,  # e.g. "Rainfall chart (rainfall.png)"
            "table": table_info
        })

    slides_data.sort(key=lambda x: x["slide_num"])

    # ----------------------------------------------------------------
    # 5) Create slides in PPTX using a BLANK layout
    #    - Title at the top center
    #    - Text in the upper portion
    #    - Chart in the lower portion if mentioned
    # ----------------------------------------------------------------
    blank_layout = prs.slide_layouts[6]  # "Blank" layout

    for sdata in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # 5.1) Add a top-centered title shape
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9.0), Inches(1.0)
        )
        title_frame = title_box.text_frame
        title_frame.text = ""
        p_title = title_frame.add_paragraph()
        p_title.text = sdata["slide_title"]
        p_title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        p_title.font.size = Pt(32)
        p_title.font.bold = True

        # 5.2) Check if there's a chart mention
        #     If so, parse out the .png filename from parentheses
        chart_filename_found = None
        if sdata["chart"]:
            match_png = re.search(r"\(([^)]+\.png)\)", sdata["chart"])
            if match_png:
                potential_name = match_png.group(1)  # e.g. "rainfall.png"
                # Check if it exists in data["chart_filename"]
                if potential_name in data["chart_filename"]:
                    chart_filename_found = potential_name

        # 5.3) If a chart is present, let's place text in the top half, chart in bottom half
        if chart_filename_found:
            # Text box in top half
            text_top = Inches(1.3)
            text_box_height = Inches(4.0)  # about half the slide
            text_left = Inches(0.5)
            text_width = Inches(7.0)

            text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_box_height)
            text_frame = text_box.text_frame
            text_frame.text = ""

            # Then the chart below that
            chart_top = Inches(3.5)  # below the text box
            chart_left = Inches(0.5)
            chart_width = Inches(7.0)
            chart_height = Inches(3.0)

        else:
            # If no chart, let the text box be larger
            text_top = Inches(1.3)
            text_box_height = Inches(4.0)
            text_left = Inches(0.5)
            text_width = Inches(7.0)

            text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_box_height)
            text_frame = text_box.text_frame
            text_frame.text = ""

            # no chart to place
            chart_top = None
            chart_left = None
            chart_width = None
            chart_height = None

        # 5.4) Fill in the content text + bullets
        if sdata["content"]:
            content_p = text_frame.add_paragraph()
            content_p.text = sdata["content"]
            content_p.level = 0
            content_p.font.size = Pt(16)

        for bullet in sdata["bullets"]:
            bp = text_frame.add_paragraph()
            bp.text = bullet
            bp.level = 1
            bp.font.size = Pt(14)

        # 5.5) If we found a matching chart file, place it below the text
        if chart_filename_found and chart_top is not None:
            if os.path.exists(chart_filename_found):
                slide.shapes.add_picture(
                    chart_filename_found,
                    chart_left,
                    chart_top,
                    chart_width,
                    chart_height
                )
            else:
                print(f"Chart file not found: {chart_filename_found}")

    # ----------------------------------------------------------------
    # 6) Save the PPTX
    # ----------------------------------------------------------------
    output_filename = "output.pptx"
    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}.")
