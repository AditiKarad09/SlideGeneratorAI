import json
import os
import re
import pandas as pd
from bs4 import BeautifulSoup
import seaborn as sns
from PyPDF2 import PdfReader
import google.generativeai as genai
import matplotlib.pyplot as plt
from processing_pdf import process_pdf  # Assuming you have a process_pdf function
import pptx
from pptx.util import Inches
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

# --------------------------------------------------------------------
# 0. CONFIGURE GEMINI
# --------------------------------------------------------------------
genai.configure(api_key="AIzaSyBq6S9KcXYfNzFiv32X0wtBen7GSEQu7gU")  # Replace with your actual key
model = genai.GenerativeModel("gemini-1.5-flash")


# --------------------------------------------------------------------
# 1A. LATEX TABLE PARSER
# --------------------------------------------------------------------
def parse_latex_table(latex_content):
    """
    Extracts the first tabular environment from the LaTeX content
    and returns a Pandas DataFrame if found. Otherwise returns None.
    """
    # 1) Regex to find first "\begin{tabular}{...} ... \end{tabular}"
    table_pattern = r"\\begin\{tabular\}\{[^\}]*\}(.*?)\\end\{tabular\}"
    match = re.search(table_pattern, latex_content, flags=re.DOTALL)
    if not match:
        # No table found in the LaTeX content
        return None

    table_body = match.group(1)

    # 2) Clean up lines (ignore \hline and empty)
    lines = []
    for line in table_body.splitlines():
        line = line.strip()
        if not line or r"\hline" in line:
            continue
        lines.append(line)

    # 3) Parse each line: remove trailing "\\" and split on "&"
    rows = []
    for line in lines:
        line = line.replace("\\\\", "").strip()
        cols = [col.strip() for col in line.split("&")]
        rows.append(cols)

    # If there's no data at all, return None
    if not rows:
        return None

    # 4) Assume the first row is the header row
    header = rows[0]
    data_rows = rows[1:]

    # 5) Convert to a DataFrame
    df = pd.DataFrame(data_rows, columns=header)
    return df


# --------------------------------------------------------------------
# 1B. FILE PROCESSING FUNCTIONS
# --------------------------------------------------------------------
def process_text(text):
    """Convert plain text into structured JSON."""
    return {"type": "text", "content": text.strip()}


def process_json(json_str):
    """Parse JSON string and return a dictionary."""
    try:
        return json.loads(json_str)
    except json.JSONDecodeError:
        return {"error": "Invalid JSON format"}


def process_table(file_path):
    """Read tabular data (CSV/Excel) and convert it to JSON."""
    try:
        if file_path.endswith(".csv"):
            df = pd.read_csv(file_path)
        elif file_path.endswith(".xlsx"):
            df = pd.read_excel(file_path)
        else:
            return {"error": "Unsupported table format"}
        return {"type": "table", "content": df.to_dict(orient="records")}
    except Exception as e:
        return {"error": f"Failed to process table: {str(e)}"}


def process_css(css_str):
    """Parse CSS and return a JSON representation."""
    try:
        soup = BeautifulSoup(css_str, "html.parser")
        rules = []
        for rule in soup.text.split("}"):
            if "{" in rule:
                selector, properties = rule.split("{", 1)
                rules.append({"selector": selector.strip(), "properties": properties.strip()})
        return {"type": "css", "content": rules}
    except Exception as e:
        return {"error": f"Failed to process CSS: {str(e)}"}


def validate_and_convert(file_path):
    """Validate and convert different file types to a dictionary representation."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return process_text(f.read())

        elif ext == ".json":
            with open(file_path, "r", encoding="utf-8") as f:
                return process_json(f.read())

        elif ext in [".csv", ".xlsx"]:
            return process_table(file_path)

        elif ext == ".css":
            with open(file_path, "r", encoding="utf-8") as f:
                return process_css(f.read())

        elif ext == ".pdf":
            # Suppose process_pdf(file_path) returns:
            #   {"type": "latex", "content": "..."} when a LaTeX doc is extracted
            pdf_result = process_pdf(file_path)
            if pdf_result.get("type") == "latex":
                latex_str = pdf_result["content"]
                # Try parsing the tabular into a DF:
                df = parse_latex_table(latex_str)
                if df is not None and not df.empty:
                    # Return table type
                    return {
                        "type": "table",
                        "content": df.to_dict(orient="records")
                    }
                else:
                    # If no table found or it's empty, just return the latex content
                    return pdf_result
            else:
                return pdf_result
        else:
            return {"error": "Unsupported file format"}
    except Exception as e:
        return {"error": f"Failed to process file: {str(e)}"}


# --------------------------------------------------------------------
# 2. CHART FUNCTIONS
# --------------------------------------------------------------------

def gemini_recommend_chart_type(table_df, purpose="visualize data"):
    """
    Calls Gemini to recommend a chart type for the given table.
    """
    columns_info = ", ".join([f"'{col}'" for col in table_df.columns])
    sample_records = table_df.head(3).to_dict(orient="records")

    prompt = f"""
I have a table with columns: {columns_info}.
Sample rows: {sample_records}.
I want to {purpose}. Which chart type is best (e.g., bar, line, scatter, pie etc.)?
Provide a short and concise answer.
"""
    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        return f"Error calling Gemini API: {str(e)}"


def plot_data_automatically(table_df, chart_type_suggestion, output_name="chart.png"):
    """
    Generates and saves a chart based on the table data and chart type suggestion.
    """
    numeric_cols = table_df.select_dtypes(include=["int", "float"]).columns
    non_numeric_cols = table_df.select_dtypes(exclude=["int", "float"]).columns
    suggestion_lower = chart_type_suggestion.lower()

    plt.figure()
    if suggestion_lower in ["line", "scatter"] and len(numeric_cols) >= 2:
        if suggestion_lower == "line":
            plt.plot(table_df[numeric_cols[0]], table_df[numeric_cols[1]], marker="o")
        else:
            plt.scatter(table_df[numeric_cols[0]], table_df[numeric_cols[1]], c="blue")
        plt.title(f"{chart_type_suggestion.capitalize()} Chart")

    elif suggestion_lower in ["hist", "distribution"] and len(numeric_cols) >= 1:
        sns.histplot(table_df[numeric_cols[0]])
        plt.title(f"{chart_type_suggestion.capitalize()} Plot")

    elif suggestion_lower in ["bar", "column"] and len(numeric_cols) >= 1 and len(non_numeric_cols) >= 1:
        plt.bar(table_df[non_numeric_cols[0]], table_df[numeric_cols[0]], color="orange")
        plt.title(f"{chart_type_suggestion.capitalize()} Chart")

    elif suggestion_lower in ["bar", "column"] and len(non_numeric_cols) >= 1:
        sns.countplot(x=table_df[non_numeric_cols[0]])
        plt.title(f"{chart_type_suggestion.capitalize()} Chart")

    else:
        print(f"Unsupported or ambiguous chart type: {chart_type_suggestion}")
        plt.close()
        return

    plt.tight_layout()
    plt.savefig(output_name)
    plt.close()


# --------------------------------------------------------------------
# 3. HELPER FOR SPLITTING BULLETS ON PERIOD
# --------------------------------------------------------------------

def split_bullets_on_period(bullet_text):
    """
    Splits bullet_text on a period + optional space => distinct bullet lines.
    """
    bullet_text = bullet_text.strip()
    lines = re.split(r"\.\s+|\.$", bullet_text)
    lines = [l.strip() for l in lines if l.strip()]
    return lines


# --------------------------------------------------------------------
# 4. HELPERS FOR REFERENCES SLIDE CLEANUP
# --------------------------------------------------------------------

def cleanup_unwanted_references_text(text):
    """
    Removes placeholder or extra lines from the references slide content.
    For example, remove lines like:
      "Remember to replace placeholder text with actual data..."
    Adjust the patterns as needed to remove more unwanted lines.
    """
    # Remove any mention of "Remember to replace placeholder text..." case-insensitively
    text = re.sub(r"(?i)remember to replace placeholder.*", "", text).strip()
    # Remove mention of "The LaTeX document provides ..." if present
    text = re.sub(r"(?i)the latex document provides.*", "", text).strip()
    return text


# --------------------------------------------------------------------
# 5. PPT GENERATION FUNCTION (UPDATED)
# --------------------------------------------------------------------
def generate_pptx_from_gemini(data, slides_title="Presentation"):
    """
    1) Calls Gemini for a proposed slide structure in the new instructions style:
       Slide 1: Introduction
       Content (no bullet points): ...
       Bullet points:
       - ...
       - ...
    2) Parses by searching for "Slide X: Title" blocks, then reading lines for content,
       bullet points, and optional Chart/Table mentions.
    3) Don't give any extra '**' for any text before placing it in PPT.
    4) Removes placeholder references if it's a References slide.
    5) Saves the final PPTX.
    """
    prs = pptx.Presentation()
    print("here ----------------------------------------------------")
    print(data)
    # --- Create Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = slides_title

    # --- 1. Build Gemini Prompt ---
    gemini_prompt = f"""
Create a PowerPoint presentation based on the following information:

**Data:** {json.dumps(data, indent=4)}

**Instructions:**
- **Structure:** Provide a clear and concise structure, specifying the content for each slide. Dont give unecessary "*"
- **Slide Titles:** Use colons (:) to separate slide numbers and titles (e.g., "Slide 1: Introduction").
- **Text Formatting:** 
   - If you have a paragraph of text, write: 
       Content (no bullet points): <text here>
   - If you have bullet points, write:
       Bullet points:
       - ...
       - ...
   -Very impornt - do not tell me to take text from anywhere just give what will be going in each slide i dont
    want to copy anything from the pdf u give what all is to be added here itself         
- **Visualizations:** If needed, mention them with lines like "- Chart:" or "- Table: ...".
- **Example**:
  Slide 1: Introduction
  Content (no bullet points): This is the opening paragraph.

  Slide 2: Key Findings
  Bullet points:
  - Finding 1
  - Finding 2

Break down information into multiple slides if it improves clarity.
"""

    try:
        response = model.generate_content(gemini_prompt)
        gemini_text = response.text.strip()
        print("[Gemini Slide Structure Suggestion]:")
        print(gemini_text)
    except Exception as e:
        print(f"Error calling Gemini API: {str(e)}")
        return

    # --- 2. Parse for "Slide X: Title" blocks ---
    slide_block_pattern = r"(Slide\s+(\d+)\s*:\s*(.*?))(?=Slide\s+\d+|$)"
    raw_blocks = re.findall(slide_block_pattern, gemini_text, flags=re.DOTALL)

    slides_data = []
    for full_block, slide_num_str, after_colon_title in raw_blocks:
        # Convert slide_num_str to int if possible
        try:
            slide_num = int(slide_num_str.strip())
        except ValueError:
            slide_num = 999  # fallback if parsing fails

        lines = full_block.splitlines()
        block_lines = lines[1:]  # everything after the "Slide X: Title" line

        # Slide title
        slide_title = after_colon_title.split("\n")[0].strip()

        content_str = ""
        bullet_points = []
        chart_info = ""
        table_info = ""
        in_bullet_mode = False

        for line in block_lines:
            txt = line.strip()
            if not txt:
                in_bullet_mode = False
                continue

            content_match = re.match(
                r"^Content\s*\(no\s*bullet\s*points\)\s*:\s*(.*)$", txt, re.IGNORECASE
            )
            bulletpts_match = re.match(r"^Bullet\s*points\s*:\s*$", txt, re.IGNORECASE)
            chart_match = re.match(r"^-?\s*Chart:\s*(.*)$", txt, re.IGNORECASE)
            table_match = re.match(r"^-?\s*Table:\s*(.*)$", txt, re.IGNORECASE)

            if content_match:
                content_str += content_match.group(1).strip() + " "
                in_bullet_mode = False
            elif bulletpts_match:
                in_bullet_mode = True
            elif chart_match:
                chart_info = chart_match.group(1).strip()
                in_bullet_mode = False
            elif table_match:
                table_info = table_match.group(1).strip()
                in_bullet_mode = False
            else:
                if in_bullet_mode and txt.startswith("-"):
                    bullet_points.append(txt[1:].strip())
                else:
                    content_str += txt + " "

        final_content = content_str.strip()

        # If the slide title is "References", remove placeholders
        if re.search(r"(?i)references", slide_title):
            final_content = cleanup_unwanted_references_text(final_content)

        # Remove "**" if present
        slide_title = slide_title.replace("**", "")
        final_content = final_content.replace("**", "")
        bullet_points = [bp.replace("**", "") for bp in bullet_points]
        chart_info = chart_info.replace("**", "")
        table_info = table_info.replace("**", "")

        slides_data.append({
            "slide_num": slide_num,
            "slide_title": slide_title,
            "content": final_content,
            "bullets": bullet_points,
            "chart": chart_info,
            "table": table_info
        })

    # Sort slides data
    slides_data.sort(key=lambda x: x["slide_num"])

    # --- 3. Build slides in PPT ---
    for sdata in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title & Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = sdata["slide_title"]

        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        # Main content
        if sdata["content"]:
            p = text_frame.add_paragraph()
            p.text = sdata["content"]
            p.level = 0

        # Bullet points
        for bullet in sdata["bullets"]:
            p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 1

        # If there's a chart mention and we have a chart file
        if sdata["chart"] and "chart_filename" in data:
            chart_path = data["chart_filename"]
            if os.path.exists(chart_path):
                left = Inches(1)
                top = Inches(2)
                width = Inches(6)
                height = Inches(4)
                slide.shapes.add_picture(chart_path, left, top, width, height)
            else:
                print("Chart file not found:", chart_path)

    # --- 4. Save the PPTX ---
    output_filename = "output2.pptx"
    prs.save(output_filename)
    print(f"Presentation saved as {output_filename}.")


# --------------------------------------------------------------------
# 6. MAIN FLOW
# --------------------------------------------------------------------
if __name__ == "__main__":
    file_path = input("Enter the file path: ")
    data = validate_and_convert(file_path)

    if "error" in data:
        print("Error processing file:", data["error"])
        exit()

    # If the data is a table, generate a chart if possible
    if data.get("type") == "table":
        df = pd.DataFrame(data["content"])
        if df.empty:
            print("No table data found.")
        else:
            chart_type = gemini_recommend_chart_type(df)
            print("Gemini suggests chart type:", chart_type)
            plot_data_automatically(df, chart_type, "chart_0.png")
            data["chart_filename"] = "chart_0.png"
    else:
        print("No table data found.")

    # Generate the PPT
    generate_pptx_from_gemini(data, slides_title="My Gemini Presentation")
